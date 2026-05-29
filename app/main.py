from fastapi import FastAPI, UploadFile, File, BackgroundTasks, Form

from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse

from fastapi.staticfiles import StaticFiles

from fastapi.middleware.cors import CORSMiddleware

from pathlib import Path

import shutil

import uuid

import datetime

import fitz
fitz.TOOLS.mupdf_display_errors(False)

import os

import asyncio

import io

import re

import fastapi

from docx import Document



from app.auth import get_current_user, get_current_active_admin, create_access_token, verify_password, ACCESS_TOKEN_EXPIRE_MINUTES, get_user_from_db

from fastapi.security import OAuth2PasswordRequestForm

from typing import Optional, List

from pydantic import BaseModel

from datetime import timedelta

import sqlite3

import json

from app.rag_engine import DB_PATH

from app.rag_engine import index_pdf_async, chat_inferential, chat_inferential_stream, document_status, DOCS_DIR, TREES_DIR, save_status, cancelled_jobs



app = FastAPI(title="Inferential RAG App with Tree Index")



# CORS setup

app.add_middleware(

    CORSMiddleware,

    allow_origins=["*"],

    allow_credentials=True,

    allow_methods=["*"],

    allow_headers=["*"],

)



# Static files for frontend

BASE_DIR = Path(__file__).parent.parent

STATIC_DIR = BASE_DIR / "static"

STATIC_DIR.mkdir(exist_ok=True)

app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

DOCS_DIR.mkdir(exist_ok=True)

app.mount("/docs", StaticFiles(directory=str(DOCS_DIR)), name="docs")



# Branding assets

BRANDING_FILE = BASE_DIR / "branding.json"

LOGOS_DIR = BASE_DIR / "logos"

LOGOS_DIR.mkdir(exist_ok=True)

app.mount("/logos", StaticFiles(directory=str(LOGOS_DIR)), name="logos")

from fastapi.responses import FileResponse

@app.get("/sw.js")
async def get_sw():
    return FileResponse(str(STATIC_DIR / "sw.js"), media_type="application/javascript")

@app.get("/manifest.json")
async def get_manifest():
    return FileResponse(str(STATIC_DIR / "manifest.json"), media_type="application/manifest+json")


# Profiles assets

PROFILES_DIR = BASE_DIR / "profiles"
PROFILES_DIR.mkdir(exist_ok=True)
app.mount("/profiles", StaticFiles(directory=str(PROFILES_DIR)), name="profiles")

# Support assets
SUPPORT_DIR = BASE_DIR / "support"
SUPPORT_DIR.mkdir(exist_ok=True)
app.mount("/support", StaticFiles(directory=str(SUPPORT_DIR)), name="support")



def load_branding() -> dict:

    if BRANDING_FILE.exists():

        try:

            with open(BRANDING_FILE, "r", encoding="utf-8") as f:

                return json.load(f)

        except Exception:

            pass

    return {"company_name": "추론형 RAG", "tagline": "", "logo_url": ""}



def save_branding(data: dict):

    with open(BRANDING_FILE, "w", encoding="utf-8") as f:

        json.dump(data, f, ensure_ascii=False, indent=2)



from app.mcp_engine import mcp_manager
import asyncio
from datetime import datetime, timedelta

SYSTEM_NOTIFICATIONS = []

async def auto_crawl_scheduler():
    import sqlite3
    from app.rag_engine import document_status, save_status, DB_PATH
    global SYSTEM_NOTIFICATIONS
    
    interval_map = {
        "5mins": timedelta(minutes=5),
        "hourly": timedelta(hours=1),
        "6hours": timedelta(hours=6),
        "12hours": timedelta(hours=12),
        "daily": timedelta(days=1),
        "weekly": timedelta(weeks=1)
    }
    
    await asyncio.sleep(10) # Delay start
    
    while True:
        try:
            now = datetime.now()
            
            with open("auto_crawl_debug.log", "a", encoding="utf-8") as debug_file:
                debug_file.write(f"\n--- {now} ---\n")
                
                # Create a snapshot to iterate safely over dynamic dict
                docs_snapshot = list(document_status.items())
                debug_file.write(f"Snapshot size: {len(docs_snapshot)}\n")
                
                for doc_id, info in docs_snapshot:
                    # 1. Ensure scheduled
                    schedule = info.get("auto_crawl_schedule", "disable")
                    if schedule == "disable" or schedule not in interval_map:
                        continue
                        
                    debug_file.write(f"Evaluating {doc_id} schedule: {schedule}\n")
                    # 2. Check if already running
                    status = info.get("status", "")
                    progress = info.get("progress", "")
                    if status in ["crawling", "processing", "pending"] or "수집 중" in progress or "진행 중" in progress:
                        debug_file.write(f"Blocked by status/progress for {doc_id}\n")
                        continue
                        
                    # 3. Check execution time matching
                    last_crawl_str = info.get("last_auto_crawl_at", "2000-01-01 00:00:00")
                    debug_file.write(f"Last crawl str: {last_crawl_str} for {doc_id}\n")
                    try:
                        last_crawl = datetime.strptime(last_crawl_str, "%Y-%m-%d %H:%M:%S")
                    except Exception:
                        last_crawl = datetime.min
                        
                    delta = interval_map.get(schedule)
                    if now - last_crawl >= delta:
                        cat_name = info.get("name", "Unknown").replace("[WEBSITE] ", "").strip()
                        debug_file.write(f"[AutoCrawl] WILL TRIGGER: {cat_name} (ID: {doc_id})\n")
                        
                        try:
                            import urllib.parse
                            import sys
                            import time
                            
                            # Use the stored file_path which contains the base url
                            url = info.get("file_path", "").replace("[WEBSITE] ", "").strip()
                            if not url.startswith("http"):
                                # Fallback inference
                                conn2 = sqlite3.connect(str(DB_PATH), timeout=30)
                                cur2 = conn2.cursor()
                                cur2.execute("SELECT url FROM web_crawl_cache WHERE doc_id = ? LIMIT 1", (doc_id,))
                                url_row = cur2.fetchone()
                                if url_row:
                                    url = url_row[0]
                                    url = '{uri.scheme}://{uri.netloc}'.format(uri=urllib.parse.urlparse(url)) if 'urllib' in sys.modules else url
                                conn2.close()
                            
                            # Prepare payload FIRST before starting thread to prevent double-scheduling due to race condition
                            document_status[doc_id]["last_auto_crawl_at"] = now.strftime("%Y-%m-%d %H:%M:%S")
                            document_status[doc_id]["status"] = "pending"
                            document_status[doc_id]["progress"] = "자동 수집 준비 중..."
                            save_status()
                            
                            if url and url.startswith("http"):
                                opts = info.get("crawl_options", {})
                                req = CrawlRequest(
                                    url=url, 
                                    site_name=cat_name, 
                                    folder_id=None, 
                                    max_depth=opts.get("max_depth", 1), 
                                    max_pages=opts.get("max_pages", 50),
                                    crawl_type=opts.get("crawl_type", "spa"),
                                    strategy=opts.get("strategy", "bfs"),
                                    restrict_path=opts.get("restrict_path", False),
                                    use_ai_extraction=opts.get("use_ai_extraction", False),
                                    ai_extraction_prompt=opts.get("ai_extraction_prompt", ""),
                                    clear_existing=False
                                )
                                owner_id = info.get("owner_id", 1)
                                org_id = info.get("organization_id", None)
                                visibility = info.get("visibility", "organization")
                                
                                asyncio.create_task(asyncio.to_thread(
                                    background_crawl_website, 
                                    doc_id, req, owner_id, org_id, cat_name, visibility
                                ))
                                SYSTEM_NOTIFICATIONS.append({
                                    "message": f"[자동업데이트] '{cat_name}'의 최신 데이터를 가져옵니다.", 
                                    "type": "info", 
                                    "time": time.time()
                                })
                        except Exception as e:
                            print(f"Error starting auto crawl for doc {doc_id}: {e}")
                            
        except Exception as e:
            print(f"[AutoCrawl] Scheduler Error: {e}")
            
        await asyncio.sleep(60) # check every minute

@app.on_event("startup")
async def startup_event():

    import json

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id, agent_type, config, python_code FROM chat_agents WHERE is_active = 1")

        rows = cursor.fetchall()

        conn.close()

        for r in rows:

            agent_id, a_type, config_text, code_text = r

            if config_text:

                try:

                    c = json.loads(config_text)

                    if a_type == 'MCP_SERVER':

                        cmd = c.get("command")

                        if cmd:

                            asyncio.create_task(mcp_manager.start_mcp_server(agent_id, cmd, c.get("args", []), c.get("env", {})))

                    elif a_type == 'MCP_CLIENT':

                        url = c.get("url")

                        if url:

                            asyncio.create_task(mcp_manager.start_mcp_client(agent_id, url, c.get("auth_token")))

                    elif a_type == 'AUTONOMOUS' and c.get('cfg-auto-daemon') == True:

                        from app.daemon_engine import daemon_manager

                        if code_text:

                            daemon_manager.start_daemon(agent_id, code_text)

                except Exception as e:

                    print(f"Failed to start Agent {agent_id} on startup: {e}")

    except Exception as e:
        print(f"Engine startup error: {e}")
        
    asyncio.create_task(auto_crawl_scheduler())

@app.get("/api/client/notifications")
async def get_system_notifications():
    """Fetches and clears memory system notifications."""
    global SYSTEM_NOTIFICATIONS
    notes = list(SYSTEM_NOTIFICATIONS)
    SYSTEM_NOTIFICATIONS.clear()
    return {"notifications": notes}

@app.on_event("shutdown")

async def shutdown_event():

    try:

        from app.mcp_engine import mcp_manager

        from app.daemon_engine import daemon_manager

        await mcp_manager.stop_all()

        daemon_manager.stop_all()

    except Exception as e:

        print(f"Engine shutdown error: {e}")





@app.get("/", response_class=HTMLResponse)

async def read_index():

    with open(STATIC_DIR / "index.html", "r", encoding="utf-8") as f:

        content = f.read()

    return HTMLResponse(content=content, headers={

        "Cache-Control": "no-store, no-cache, must-revalidate",

        "Pragma": "no-cache",

    })



@app.get("/js/script.js")

async def serve_script():

    """Serve script.js with no-cache headers to prevent browser caching issues."""

    from fastapi.responses import FileResponse

    return FileResponse(

        path=str(STATIC_DIR / "script.js"),

        media_type="application/javascript",

        headers={

            "Cache-Control": "no-store, no-cache, must-revalidate",

            "Pragma": "no-cache",

        }

    )



class LoginRequest(BaseModel):

    username: str

    password: str



@app.post("/api/login")

async def login_for_access_token(form_data: OAuth2PasswordRequestForm = fastapi.Depends()):

    user = get_user_from_db(form_data.username)

    if not user or not verify_password(form_data.password, user["password_hash"]):

        return JSONResponse(status_code=401, content={"error": "Incorrect username or password"})

    

    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)

    access_token = create_access_token(

        data={"sub": user["username"], "role": user["role"]}, expires_delta=access_token_expires

    )

    return {

        "access_token": access_token, 

        "token_type": "bearer", 

        "role": user["role"], 

        "username": user["username"], 

        "id": user["id"], 

        "organization_id": user["organization_id"],

        "profile_image": user.get("profile_image")

    }



class PasswordChangeRequest(BaseModel):

    current_password: str

    new_password: str



@app.get("/api/users/me")

async def read_users_me(current_user: dict = fastapi.Depends(get_current_user)):

    return {

        "id": current_user["id"],

        "username": current_user["username"],

        "role": current_user["role"],

        "organization_id": current_user["organization_id"],

        "profile_image": current_user.get("profile_image")

    }



@app.post("/api/users/me/password")

async def change_password(request: PasswordChangeRequest, current_user: dict = fastapi.Depends(get_current_user)):

    user_db = get_user_from_db(current_user["username"])

    if not verify_password(request.current_password, user_db["password_hash"]):

        return JSONResponse(status_code=400, content={"error": "현재 비밀번호가 일치하지 않습니다."})

    

    from app.auth import get_password_hash

    hashed_pw = get_password_hash(request.new_password)

    

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    cursor = conn.cursor()

    cursor.execute("UPDATE users SET password_hash = ? WHERE id = ?", (hashed_pw, current_user["id"]))

    conn.commit()

    conn.close()

    

    return {"message": "ë¹„ë?ë²ˆí˜¸ê°€ ?±ê³µ?�ìœ¼ë¡?ë³€ê²½ë�˜?ˆìŠµ?ˆë‹¤."}



@app.post("/api/users/me/profile_image")

async def upload_profile_image(file: UploadFile = File(...), current_user: dict = fastapi.Depends(get_current_user)):

    allowed_exts = ['.png', '.jpg', '.jpeg', '.gif', '.webp']

    ext = os.path.splitext(file.filename)[1].lower()

    if ext not in allowed_exts:

        return JSONResponse(status_code=400, content={"error": f"ì§€?�í•˜ì§€ ?ŠëŠ” ?´ë?ì§€ ?•ì‹�?…ë‹ˆ?? {ext}"})

        

    safe_filename = f"user_{current_user['id']}{ext}"

    file_path = PROFILES_DIR / safe_filename

    

    with open(file_path, "wb") as buffer:

        shutil.copyfileobj(file.file, buffer)

        

    # Append a timestamp to avoid browser caching issues when updating the image

    profile_url = f"/profiles/{safe_filename}?v={int(datetime.now().timestamp())}"

    

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    cursor = conn.cursor()

    cursor.execute("UPDATE users SET profile_image = ? WHERE id = ?", (profile_url, current_user["id"]))

    conn.commit()

    conn.close()

    

    return {"message": "?„ë¡œ???´ë?ì§€ê°€ ?±ê³µ?�ìœ¼ë¡?ë³€ê²½ë�˜?ˆìŠµ?ˆë‹¤.", "profile_image": profile_url}



@app.get("/api/admin/stats/usage")

async def get_usage_stats(start_date: str = None, end_date: str = None, current_user: dict = fastapi.Depends(get_current_user)):

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    conn.row_factory = sqlite3.Row

    cursor = conn.cursor()

    

    # 1. Total tokens per user (for ranking table)

    query_user = """

        SELECT u.id, u.username, u.full_name, t.model_name,

               SUM(t.prompt_tokens) as total_prompt, 

               SUM(t.completion_tokens) as total_completion,

               MAX(p.cost_per_1m_prompt) as cost_prompt,

               MAX(p.cost_per_1m_completion) as cost_comp

        FROM token_usage t

        JOIN users u ON t.user_id = u.id

        LEFT JOIN model_pricing p ON t.model_name = p.model_name

        WHERE 1=1

    """

    params = []
    
    if current_user.get("role") != "admin":
        query_user += " AND t.user_id = ?"
        params.append(current_user["id"])

    if start_date:

        query_user += " AND t.created_at >= ?"

        params.append(start_date + " 00:00:00")

    if end_date:

        query_user += " AND t.created_at <= ?"

        params.append(end_date + " 23:59:59")

        

    query_user += " GROUP BY u.id, u.username, u.full_name, t.model_name"

    

    cursor.execute(query_user, params)

    user_rows = cursor.fetchall()

    

    # Process user rankings and calculate costs

    user_stats = {}

    for r in user_rows:

        uid = r["id"]

        if uid not in user_stats:

            user_stats[uid] = {

                "id": uid,

                "username": r["username"],

                "full_name": r["full_name"],

                "total_tokens": 0,

                "total_prompt": 0,

                "total_completion": 0,

                "total_cost": 0.0,

                "models": {}

            }

            

        pt = r["total_prompt"] or 0

        ct = r["total_completion"] or 0

        cp = r["cost_prompt"] or 0.0
        cc = r["cost_comp"] or 0.0
        
        model_name = r["model_name"]
        if cp == 0.0:
            if model_name == "gemini-flash-lite-latest": cp = 0.75
            elif model_name == "gemini-flash-latest": cp = 1.50
            
        if cc == 0.0:
            if model_name == "gemini-flash-lite-latest": cc = 4.50
            elif model_name == "gemini-flash-latest": cc = 9.00

        

        cost = ((pt / 1000000.0) * cp) + ((ct / 1000000.0) * cc)

        

        user_stats[uid]["total_prompt"] += pt

        user_stats[uid]["total_completion"] += ct

        user_stats[uid]["total_tokens"] += (pt + ct)

        user_stats[uid]["total_cost"] += cost

        user_stats[uid]["models"][r["model_name"]] = {

            "prompt": pt, "completion": ct, "cost": cost

        }

        

    # Sort users by total custom (descending)

    user_rankings = sorted(list(user_stats.values()), key=lambda x: x["total_tokens"], reverse=True)

        

    # 2. Daily trends for the chart

    query_daily = """

        SELECT date(t.created_at) as usage_date, t.model_name,

               SUM(t.prompt_tokens + t.completion_tokens) as daily_tokens

        FROM token_usage t

        WHERE 1=1

    """

    params_daily = []
    
    if current_user.get("role") != "admin":
        query_daily += " AND t.user_id = ?"
        params_daily.append(current_user["id"])

    if start_date:

        query_daily += " AND t.created_at >= ?"

        params_daily.append(start_date + " 00:00:00")

    if end_date:

        query_daily += " AND t.created_at <= ?"

        params_daily.append(end_date + " 23:59:59")

        

    query_daily += " GROUP BY usage_date, t.model_name ORDER BY usage_date ASC"

    

    cursor.execute(query_daily, params_daily)

    daily_rows = cursor.fetchall()

    

    daily_stats = {}

    for r in daily_rows:

        d = r["usage_date"]

        m = r["model_name"]

        tk = r["daily_tokens"]

        if d not in daily_stats:

            daily_stats[d] = {}

        daily_stats[d][m] = tk

        

    conn.close()

    

    return {

        "user_rankings": user_rankings,

        "daily_trends": daily_stats

    }



class PricingUpdate(BaseModel):

    model_name: str

    cost_per_1m_prompt: float

    cost_per_1m_completion: float



@app.get("/api/admin/stats/pricing")

async def get_pricing(current_user: dict = fastapi.Depends(get_current_active_admin)):

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    conn.row_factory = sqlite3.Row

    cursor = conn.cursor()

    cursor.execute("SELECT * FROM model_pricing")

    rows = cursor.fetchall()

    conn.close()

    return [dict(r) for r in rows]



@app.post("/api/admin/stats/pricing")

async def update_pricing(data: PricingUpdate, current_user: dict = fastapi.Depends(get_current_active_admin)):

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    cursor = conn.cursor()

    cursor.execute('''

        INSERT INTO model_pricing (model_name, cost_per_1m_prompt, cost_per_1m_completion)

        VALUES (?, ?, ?)

        ON CONFLICT(model_name) DO UPDATE SET

        cost_per_1m_prompt=excluded.cost_per_1m_prompt,

        cost_per_1m_completion=excluded.cost_per_1m_completion

    ''', (data.model_name, data.cost_per_1m_prompt, data.cost_per_1m_completion))

    conn.commit()
    conn.close()
    return {"message": "단가 설정이 갱신되었습니다."}

@app.delete("/api/admin/stats/pricing/{model_name}")

async def delete_pricing(model_name: str, current_user: dict = fastapi.Depends(get_current_active_admin)):

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    cursor = conn.cursor()

    cursor.execute("DELETE FROM model_pricing WHERE model_name=?", (model_name,))

    conn.commit()
    conn.close()
    return {"message": "해당 단가 설정이 삭제되었습니다."}

@app.post("/api/external/upload")
async def external_upload_pdf(
    background_tasks: BackgroundTasks, 
    file: UploadFile = File(...), 
    category: str = Form("General"), 
    visibility: str = Form("organization"),
    owner_username: str = Form(None),
    x_api_key: str = fastapi.Header(None)
):
    expected_key = os.environ.get("M2M_UPLOAD_API_KEY", "aimm-server-api-key")
    if not x_api_key or x_api_key != expected_key:
        return JSONResponse(status_code=401, content={"error": "Invalid API Key"})
    
    # Check for duplicates based on name and category (robust check)
    target_category = (category or "General").strip()
    target_filename = (file.filename or "").strip()
    for val in document_status.values():
        if not val.get("is_active", True) or val.get("status") == "error":
            continue
        val_cat = (val.get("category") or "General").strip()
        val_name = (val.get("name") or "").strip()
        if val_name == target_filename and val_cat == target_category:
            return JSONResponse(status_code=409, content={"error": "동일한 폴더에 같은 이름의 문서가 이미 존재합니다.", "duplicate": True})
    


    conn = sqlite3.connect(str(DB_PATH), timeout=30)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    
    if owner_username:
        cursor.execute("SELECT * FROM users WHERE username = ? LIMIT 1", (owner_username,))
        admin_user = cursor.fetchone()
        if not admin_user:
            conn.close()
            return JSONResponse(status_code=400, content={"error": f"지정한 사용자({owner_username})를 찾을 수 없습니다."})
    else:
        cursor.execute("SELECT * FROM users WHERE role='admin' ORDER BY id LIMIT 1")
        admin_user = cursor.fetchone()
        
    if not admin_user:
        conn.close()
        return JSONResponse(status_code=500, content={"error": "시스템에 문서를 소유할 계정이 없습니다."})
        
    current_user = dict(admin_user)
    
    if category and category.strip() and category.strip() != "General":
        clean_cat = category.strip()
        cursor.execute("SELECT id FROM categories WHERE name = ? LIMIT 1", (clean_cat,))
        cat_exists = cursor.fetchone()
        if not cat_exists:
            cat_vis = visibility if visibility else 'private'
            cursor.execute("INSERT INTO categories (name, visibility, owner_id) VALUES (?, ?, ?)", (clean_cat, cat_vis, current_user["id"]))
            conn.commit()
            
    conn.close()
    
    allowed_exts = ['.pdf', '.txt', '.md', '.csv', '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.hwp', '.hwpx']
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in allowed_exts:
        return JSONResponse(status_code=400, content={"error": f"지원하지 않는 파일 형식입니다. {ext}"})
        
    doc_id = str(uuid.uuid4())[:8]
    safe_filename = f"{doc_id}_{file.filename}"
    file_path = DOCS_DIR / safe_filename
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    upload_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # Async Remote Conversion for HWP/HWPX
    if ext in ['.hwp', '.hwpx']:
        pdf_safe_filename = f"{doc_id}_{os.path.splitext(file.filename)[0]}.pdf"
        pdf_file_path = DOCS_DIR / pdf_safe_filename
        
        try:
            import requests
            with open(file_path, "rb") as f:
                resp = requests.post("https://hwp.aimm.pro:8877/convert", files={"file": f}, timeout=30)
            if resp.status_code not in (200, 202):
                raise Exception(f"Server returned {resp.status_code}: {resp.text}")
            job = resp.json()
            job_id = job["job_id"]
        except Exception as e:
            try:
                if os.path.exists(file_path): os.remove(file_path)
            except: pass
            return JSONResponse(status_code=500, content={"error": f"HWP 서버 연동 실패: {str(e)}"})
            
        document_status[doc_id] = {
            "status": "pending", 
            "name": file.filename,
            "safe_filename": safe_filename, # point to hwp until replaced
            "progress": f"PDF 변환대기 (큐: {job.get('queue_position')})",
            "upload_date": upload_date,
            "page_count": "Unknown",
            "category": category,
            "visibility": visibility,
            "is_active": True,
            "owner_id": current_user["id"],
            "organization_id": current_user["organization_id"]
        }
        save_status()
        background_tasks.add_task(poll_hwp_convert_and_index_async, doc_id, job_id, str(file_path), str(pdf_file_path))
        return {"doc_id": doc_id, "message": "HWP 업로드 및 변환 작업이 시작되었습니다.", "filename": file.filename, "safe_filename": safe_filename}

    # Convert to PDF if necessary (Local sync fallbacks for Word, PPT, etc.)
    if ext != '.pdf':
        try:
            from app.converters import convert_to_pdf
            pdf_safe_filename = f"{doc_id}_{os.path.splitext(file.filename)[0]}.pdf"
            pdf_file_path = DOCS_DIR / pdf_safe_filename
            
            success = convert_to_pdf(str(file_path), str(pdf_file_path))
            if not success:
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                except Exception:
                    pass
                return JSONResponse(status_code=500, content={"error": "문서를 PDF로 변환하는데 실패했습니다. (서버의 LibreOffice 등 설정 확인)"})
                
            # Delete the original file, we only need the converted PDF
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception:
                pass
                
            file_path = pdf_file_path
            safe_filename = pdf_safe_filename
        except Exception as e:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception:
                pass
            import traceback
            error_details = traceback.format_exc()
            return JSONResponse(status_code=500, content={"error": f"변환 중 내부 오류 발생: {str(e)}", "details": error_details})
        
    # Extract page count
    try:
        pdf_doc = fitz.open(file_path)
        page_count = len(pdf_doc)
        pdf_doc.close()
    except Exception:
        page_count = "Unknown"
        
    upload_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
    # Start indexing in background
    document_status[doc_id] = {
        "status": "pending", 
        "name": file.filename,  # Keep original name for display
        "safe_filename": safe_filename, 
        "progress": "Waiting...",
        "upload_date": upload_date,
        "page_count": page_count,
        "category": category,
        "visibility": visibility,
        "is_active": True,
        "owner_id": current_user["id"],
        "organization_id": current_user["organization_id"]
    }
    save_status()
    
    # Execute the purely non-blocking wrapped workflow in the background
    background_tasks.add_task(index_pdf_async, str(file_path), doc_id)
    
    return {"doc_id": doc_id, "message": "File uploaded, converted, and indexing started", "filename": file.filename, "safe_filename": safe_filename}

async def poll_hwp_convert_and_index_async(doc_id: str, job_id: str, hwp_file_path: str, pdf_file_path: str):
    import asyncio
    import requests
    import os
    import fitz
    
    while True:
        try:
            status_resp = requests.get(f"https://hwp.aimm.pro:8877/status/{job_id}", timeout=10)
            if status_resp.status_code == 200:
                info = status_resp.json()
                if info["status"] == "done":
                    document_status[doc_id]["progress"] = "PDF 변환 완료. 다운로드중..."
                    save_status()
                    break
                elif info["status"] == "failed":
                    document_status[doc_id]["status"] = "failed"
                    document_status[doc_id]["progress"] = f"HWP 변환 실패: {info.get('error_message')}"
                    save_status()
                    return
                elif info["status"] == "processing":
                    document_status[doc_id]["progress"] = "PDF 변환중..."
                    save_status()
                elif info["status"] == "in_queue":
                    document_status[doc_id]["progress"] = f"PDF 변환 대기 (큐: {info.get('queue_position')})"
                    save_status()
        except requests.exceptions.RequestException:
            pass # ignore network errors and retry
        await asyncio.sleep(3)
        
    try:
        pdf_resp = requests.get(f"https://hwp.aimm.pro:8877/download/{job_id}", timeout=60)
        if pdf_resp.status_code != 200:
            raise Exception("다운로드 실패 (상태 코드: " + str(pdf_resp.status_code) + ")")
        with open(pdf_file_path, "wb") as f:
            f.write(pdf_resp.content)
    except Exception as e:
        document_status[doc_id]["status"] = "failed"
        document_status[doc_id]["progress"] = f"HWP 다운로드 실패: {e}"
        save_status()
        return

    try:
        if os.path.exists(hwp_file_path):
            os.remove(hwp_file_path)
    except: pass

    try:
        pdf_doc = fitz.open(pdf_file_path)
        document_status[doc_id]["page_count"] = len(pdf_doc)
        pdf_doc.close()
    except Exception:
        pass

    document_status[doc_id]["safe_filename"] = os.path.basename(pdf_file_path)
    document_status[doc_id]["progress"] = "Waiting..."
    save_status()
    
    await index_pdf_async(str(pdf_file_path), doc_id)

@app.post("/api/upload")

async def upload_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), category: str = Form("General"), visibility: str = Form("private"), current_user: dict = fastapi.Depends(get_current_user)):

    # Check for duplicates based on name and category (robust check)
    target_category = (category or "General").strip()
    target_filename = (file.filename or "").strip()
    for val in document_status.values():
        if not val.get("is_active", True) or val.get("status") == "error":
            continue
        val_cat = (val.get("category") or "General").strip()
        val_name = (val.get("name") or "").strip()
        if val_name == target_filename and val_cat == target_category:
            return JSONResponse(status_code=409, content={"error": "동일한 대상 폴더에 같은 이름의 문서가 이미 존재합니다.", "duplicate": True})

    allowed_exts = ['.pdf', '.txt', '.md', '.csv', '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.hwp', '.hwpx']

    ext = os.path.splitext(file.filename)[1].lower()

    if ext not in allowed_exts:

        return JSONResponse(status_code=400, content={"error": f"ì§€?�í•˜ì§€ ?ŠëŠ” ?Œì�¼ ?•ì‹�?…ë‹ˆ?? {ext}"})

        

    doc_id = str(uuid.uuid4())[:8]

    safe_filename = f"{doc_id}_{file.filename}"

    file_path = DOCS_DIR / safe_filename

    

    with open(file_path, "wb") as buffer:

        shutil.copyfileobj(file.file, buffer)

        

    upload_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # Async Remote Conversion for HWP/HWPX
    if ext in ['.hwp', '.hwpx']:
        pdf_safe_filename = f"{doc_id}_{os.path.splitext(file.filename)[0]}.pdf"
        pdf_file_path = DOCS_DIR / pdf_safe_filename
        
        try:
            import requests
            with open(file_path, "rb") as f:
                resp = requests.post("https://hwp.aimm.pro:8877/convert", files={"file": f}, timeout=30)
            if resp.status_code not in (200, 202):
                raise Exception(f"Server returned {resp.status_code}: {resp.text}")
            job = resp.json()
            job_id = job["job_id"]
        except Exception as e:
            try:
                if os.path.exists(file_path): os.remove(file_path)
            except: pass
            return JSONResponse(status_code=500, content={"error": f"HWP 서버 연동 실패: {str(e)}"})
            
        document_status[doc_id] = {
            "status": "pending", 
            "name": file.filename,
            "safe_filename": safe_filename, # point to hwp until replaced
            "progress": f"PDF 변환대기 (큐: {job.get('queue_position')})",
            "upload_date": upload_date,
            "page_count": "Unknown",
            "category": category,
            "visibility": visibility,
            "is_active": True,
            "owner_id": current_user["id"],
            "organization_id": current_user["organization_id"]
        }
        save_status()
        background_tasks.add_task(poll_hwp_convert_and_index_async, doc_id, job_id, str(file_path), str(pdf_file_path))
        return {"doc_id": doc_id, "message": "HWP 업로드 및 변환 작업이 시작되었습니다.", "filename": file.filename, "safe_filename": safe_filename}


    # Convert to PDF if necessary (Local sync fallbacks for Word, PPT, etc.)

    if ext != '.pdf':

        try:

            from app.converters import convert_to_pdf

            pdf_safe_filename = f"{doc_id}_{os.path.splitext(file.filename)[0]}.pdf"

            pdf_file_path = DOCS_DIR / pdf_safe_filename

            

            success = convert_to_pdf(str(file_path), str(pdf_file_path))

            if not success:

                try:

                    if os.path.exists(file_path):

                        os.remove(file_path)

                except Exception:

                    pass

                return JSONResponse(status_code=500, content={"error": "ë¬¸ì„œë¥?PDFë¡?ë³€?˜í•˜?????¤íŒ¨?ˆìŠµ?ˆë‹¤. (?œë²„??LibreOffice ???¤ì • ?•ì�¸)"})

                

            # Delete the original file, we only need the converted PDF

            try:

                if os.path.exists(file_path):

                    os.remove(file_path)

            except Exception:

                pass

                

            file_path = pdf_file_path

            safe_filename = pdf_safe_filename

        except Exception as e:

            try:

                if os.path.exists(file_path):

                    os.remove(file_path)

            except Exception:

                pass

            import traceback

            error_details = traceback.format_exc()

            return JSONResponse(status_code=500, content={"error": f"ë³€??ì¤??´ë? ?¤ë¥˜ ë°œìƒ�: {str(e)}", "details": error_details})

        

    # Extract page count

    try:

        pdf_doc = fitz.open(file_path)

        page_count = len(pdf_doc)

        pdf_doc.close()

    except Exception:

        page_count = "Unknown"

        

    upload_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        

    # Start indexing in background

    document_status[doc_id] = {

        "status": "pending", 

        "name": file.filename,  # Keep original name for display

        "safe_filename": safe_filename, 

        "progress": "Waiting...",

        "upload_date": upload_date,

        "page_count": page_count,

        "category": category,

        "visibility": visibility,

        "is_active": True,

        "owner_id": current_user["id"],

        "organization_id": current_user["organization_id"]

    }

    save_status()

    

    # Execute the purely non-blocking wrapped workflow in the background

    background_tasks.add_task(index_pdf_async, str(file_path), doc_id)

    

    return {"doc_id": doc_id, "message": "File uploaded, converted, and indexing started", "filename": file.filename, "safe_filename": safe_filename}



@app.post("/api/chat/upload_temp")

async def upload_temp_file(file: UploadFile = File(...), current_user: dict = fastapi.Depends(get_current_user)):

    try:

        tmpdir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "static", "agent_inputs")

        os.makedirs(tmpdir, exist_ok=True)

        import uuid

        safe_filename = f"{uuid.uuid4()}_{file.filename}"

        filepath = os.path.join(tmpdir, safe_filename)

        with open(filepath, "wb") as f:

            f.write(await file.read())

        return {"filepath": filepath, "filename": file.filename}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/extract_text")

async def extract_text(file: UploadFile = File(...)):

    if not file.filename.lower().endswith(".pdf"):

        return JSONResponse(status_code=400, content={"error": "Only PDF files are allowed"})

        

    try:

        # Read the file into memory

        content = await file.read()

        

        # Open PDF from memory using fitz

        pdf_doc = fitz.open(stream=content, filetype="pdf")

        

        extracted_text = []

        for page_num in range(len(pdf_doc)):

            page = pdf_doc[page_num]

            page_text = page.get_text("text")

            

            # OCR Fallback

            if not page_text.strip():

                try:

                    pix = page.get_pixmap(dpi=150)

                    img_bytes = pix.tobytes("png")

                    from pageindex.utils import get_gemini_client

                    from google.genai import types

                    client = get_gemini_client()

                    prompt_text = "Please extract all the text from this document image accurately. Preserve the original language and formatting as much as possible. Do not include any explanations."

                    message_parts = [

                        types.Part.from_text(text=prompt_text),

                        types.Part.from_bytes(data=img_bytes, mime_type="image/png")

                    ]

                    response = client.models.generate_content(

                        model='gemini-flash-lite-latest',

                        contents=message_parts

                    )

                    if response and response.text:

                        page_text = response.text

                        print(f"Chat Attach OCR: Extracted {len(page_text)} chars from page {page_num+1}")

                    else:

                        page_text = "\n[?´ë?ì§€ ?�ìŠ¤???¸ì‹� ?¤íŒ¨]\n"

                except Exception as e:

                    print(f"Chat Attach OCR Error: {e}")

                    page_text = "\n[?´ë?ì§€ ?�ìŠ¤???¸ì‹� ?¤íŒ¨]\n"



            extracted_text.append(page_text)

            

        pdf_doc.close()

        

        full_text = "\n".join(extracted_text)

        

        # Limit text size if it's excessively large to protect the context window

        if len(full_text) > 30000:

            full_text = full_text[:30000] + "\n...(ë¬¸ì„œê°€ ?ˆë¬´ ê¸¸ì–´ ?¼ë?ê°€ ?�ëžµ?˜ì—ˆ?µë‹ˆ??..."

            

        return {"filename": file.filename, "text": full_text.strip()}

        

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": f"Failed to extract text: {str(e)}"})



@app.get("/api/documents")

async def get_documents(current_user: dict = fastapi.Depends(get_current_user)):

    """Returns the list of indexed and indexing documents."""

    my_docs = []

    public_docs = []

    

    # Fetch all users and orgs once for mapping

    user_map = {}

    org_map = {}

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id, username, full_name, organization_id FROM users")

        for r in cursor.fetchall():

            user_map[r[0]] = {"username": r[1], "full_name": r[2], "org_id": r[3]}

        cursor.execute("SELECT id, name FROM organizations")

        for r in cursor.fetchall():

            org_map[r[0]] = r[1]

        conn.close()

    except Exception as e:

        print(f"Error fetching users/orgs for documents: {e}")

        

    user_accessible_groups = set()

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id FROM sharing_groups WHERE owner_id = ?", (current_user["id"],))

        for r in cursor.fetchall(): user_accessible_groups.add(f"group_{r[0]}")

        cursor.execute("SELECT group_id FROM sharing_group_members WHERE target_type = 'user' AND target_id = ?", (current_user["id"],))

        for r in cursor.fetchall(): user_accessible_groups.add(f"group_{r[0]}")

        if current_user.get("organization_id"):

            cursor.execute("SELECT group_id FROM sharing_group_members WHERE target_type = 'organization' AND target_id = ?", (current_user["organization_id"],))

            for r in cursor.fetchall(): user_accessible_groups.add(f"group_{r[0]}")

        conn.close()

    except Exception as e:

        print(f"Error fetching group memberships: {e}")

        

    for doc_id, info in document_status.items():

        doc_owner_id = info.get("owner_id")

        doc_org_id = info.get("organization_id")

        visibility = info.get("visibility", "public") # legacy docs act as public

        

        uploader_name = "?????†ì�Œ"

        uploader_org = "?????†ì�Œ"

        

        if doc_owner_id and doc_owner_id in user_map:

            uploader_name = user_map[doc_owner_id]["full_name"] or user_map[doc_owner_id]["username"]

            # Always use the user's CURRENT organization (not the stale doc org_id)

            current_owner_org_id = user_map[doc_owner_id]["org_id"]

            if current_owner_org_id and current_owner_org_id in org_map:

                uploader_org = org_map[current_owner_org_id]

        elif doc_org_id and doc_org_id in org_map:

            # Fallback for legacy docs without a known owner

            uploader_org = org_map[doc_org_id]

            

        doc_obj = {

            "id": doc_id,

            "status": info.get("status"),

            "progress": info.get("progress", ""),

            "progress_percent": info.get("progress_percent", 0),

            "name": info.get("name", f"Document {doc_id}"),

            "safe_filename": info.get("safe_filename"),

            "file_path": info.get("file_path", ""),

            "upload_date": info.get("upload_date", "Unknown"),

            "updated_at": info.get("updated_at", info.get("upload_date", "Unknown")),

            "page_count": info.get("page_count", "Unknown"),

            "category": info.get("category", "General"),

            "visibility": visibility,

            "doc_description": info.get("doc_description", ""),

            "auto_crawl_schedule": info.get("auto_crawl_schedule", "disable"),
            "crawl_options": info.get("crawl_options", {}),
            "is_active": info.get("is_active", True),

            "owner_id": doc_owner_id,

            "uploader_name": uploader_name,

            "uploader_org": uploader_org

        }

        

        # Admin logic: my docs go to my_docs; other users' docs follow same visibility rules

        if current_user["role"] == "admin":

            if doc_owner_id == current_user["id"]:

                my_docs.append(doc_obj)

            elif visibility == "public":

                public_docs.append(doc_obj)

            elif visibility == "organization":

                # Admins also only see org-shared docs from their own organization

                owner_current_org = user_map.get(doc_owner_id, {}).get("org_id") if doc_owner_id else None

                viewer_org = current_user["organization_id"]

                if owner_current_org and owner_current_org == viewer_org:

                    public_docs.append(doc_obj)

            elif visibility.startswith("group_") and visibility in user_accessible_groups:

                public_docs.append(doc_obj)

            # "private" docs from others are never shown

            continue

            

        # Normal user logic

        if doc_owner_id == current_user["id"]:

            my_docs.append(doc_obj)

        elif visibility == "public":

            public_docs.append(doc_obj)

        elif visibility == "organization":

            # Compare owner's CURRENT organization vs viewer's CURRENT organization

            # This ensures org changes are reflected immediately without relying on stale doc_org_id

            owner_current_org = user_map.get(doc_owner_id, {}).get("org_id") if doc_owner_id else None

            viewer_org = current_user["organization_id"]

            if owner_current_org and owner_current_org == viewer_org:

                public_docs.append(doc_obj)

        elif visibility.startswith("group_") and visibility in user_accessible_groups:

            public_docs.append(doc_obj)

            

    my_docs.sort(key=lambda x: str(x.get("upload_date", "")), reverse=True)

    public_docs.sort(key=lambda x: str(x.get("upload_date", "")), reverse=True)

    return {"my_documents": my_docs, "public_documents": public_docs}



@app.get("/api/categories")

async def get_categories():

    """Legacy: returns flat category list for backwards compatibility."""

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    conn.row_factory = sqlite3.Row

    cursor = conn.cursor()

    cursor.execute("SELECT name FROM categories ORDER BY name")

    rows = cursor.fetchall()

    conn.close()

    categories = [r["name"] for r in rows] if rows else ["General"]

    return JSONResponse(content={"categories": categories}, media_type="application/json; charset=utf-8")



# ?€?€?€ Folder (Category) CRUD ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€



@app.get("/api/folders")

async def get_folders(current_user: dict = fastapi.Depends(get_current_user)):

    """Returns accessible folders as a flat list."""

    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    conn.row_factory = sqlite3.Row

    cursor = conn.cursor()

    

    # User can see:

    # 1. Folders they own

    # 2. Folders with visibility 'public'

    # 3. Folders with visibility 'organization' where the folder owner is in the same organization

    

    user_id = current_user["id"]

    org_id = current_user.get("organization_id")

    

    # 3.5. Folders shared via custom groups

    user_accessible_groups = []

    try:

        cursor.execute("SELECT id FROM sharing_groups WHERE owner_id = ?", (user_id,))

        for r in cursor.fetchall(): user_accessible_groups.append(f"group_{r['id']}")

        cursor.execute("SELECT group_id FROM sharing_group_members WHERE target_type = 'user' AND target_id = ?", (user_id,))

        for r in cursor.fetchall(): user_accessible_groups.append(f"group_{r['group_id']}")

        if org_id:

            cursor.execute("SELECT group_id FROM sharing_group_members WHERE target_type = 'organization' AND target_id = ?", (org_id,))

            for r in cursor.fetchall(): user_accessible_groups.append(f"group_{r['group_id']}")

    except Exception as e:

        print(f"Error fetching group memberships for folders: {e}")



    group_placeholders = ','.join(['?'] * len(user_accessible_groups))

    group_clause = f"OR c.visibility IN ({group_placeholders})" if user_accessible_groups else ""

    

    if org_id is not None:

        query = f"""

            SELECT c.id, c.name, c.parent_id, c.visibility, c.owner_id, c.auto_crawl_schedule, 
                   u.full_name, u.username, o.name as org_name

            FROM categories c

            LEFT JOIN users u ON c.owner_id = u.id

            LEFT JOIN organizations o ON u.organization_id = o.id

            WHERE c.owner_id = ? 

               OR c.visibility = 'public' 

               OR (c.visibility = 'organization' AND u.organization_id = ?)

               {group_clause}

            ORDER BY c.parent_id, c.name

        """

        params = [user_id, org_id] + user_accessible_groups

        cursor.execute(query, params)

    else:

        # User has no organization, can only see own, public, or specific groups

        query = f"""

            SELECT c.id, c.name, c.parent_id, c.visibility, c.owner_id, c.auto_crawl_schedule, 
                   u.full_name, u.username, o.name as org_name

            FROM categories c

            LEFT JOIN users u ON c.owner_id = u.id

            LEFT JOIN organizations o ON u.organization_id = o.id

            WHERE c.owner_id = ? OR c.visibility = 'public'

               {group_clause}

            ORDER BY c.parent_id, c.name

        """

        params = [user_id] + user_accessible_groups

        cursor.execute(query, params)

        

    rows = cursor.fetchall()

    conn.close()

    

    folders = []

    for r in rows:

        name_part = r["full_name"] if r["full_name"] else r["username"]
        org_part = r["org_name"] if r["org_name"] else "소속 미상"
        if name_part is None:
            owner_name = "시스템 기본 항목"
        else:
            owner_name = f"{org_part}/{name_part}"

        folders.append({
            "id": r["id"], 
            "name": r["name"], 
            "parent_id": r["parent_id"], 
            "visibility": r["visibility"], 
            "owner_id": r["owner_id"], 
            "auto_crawl_schedule": r["auto_crawl_schedule"] or "disable",
            "owner_name": owner_name
        })

    return JSONResponse(content={"folders": folders}, media_type="application/json; charset=utf-8")





@app.post("/api/folders")

async def create_folder(request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    """Create a new folder (category) with visibility."""

    name = (request.get("name") or "").strip()

    parent_id = request.get("parent_id")  # None = root

    visibility = request.get("visibility", "private") # default private
    auto_crawl_schedule = request.get("auto_crawl_schedule", "disable")

    if visibility not in ("private", "organization", "public"):
        visibility = "private"

    if not name:
        return JSONResponse(status_code=400, content={"error": "Folder name is required"}, media_type="application/json; charset=utf-8")

    conn = sqlite3.connect(str(DB_PATH), timeout=30)
    cursor = conn.cursor()

    # Validate parent exists if provided
    if parent_id is not None:
        cursor.execute("SELECT id FROM categories WHERE id = ?", (parent_id,))
        if not cursor.fetchone():
            conn.close()
            return JSONResponse(status_code=404, content={"error": "Parent folder not found"}, media_type="application/json; charset=utf-8")

    try:
        cursor.execute(
            "INSERT INTO categories (name, parent_id, owner_id, visibility, auto_crawl_schedule) VALUES (?, ?, ?, ?, ?)",
            (name, parent_id, current_user["id"], visibility, auto_crawl_schedule)
        )
        folder_id = cursor.lastrowid

        conn.commit()

    except Exception as e:

        conn.close()

        return JSONResponse(status_code=500, content={"error": str(e)}, media_type="application/json; charset=utf-8")

    conn.close()

    return JSONResponse(content={"id": folder_id, "name": name, "parent_id": parent_id, "visibility": visibility}, media_type="application/json; charset=utf-8")





@app.patch("/api/folders/{folder_id}")

async def rename_folder(folder_id: int, request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    """Rename a folder and update visibility."""

    new_name = (request.get("name") or "").strip()
    visibility = request.get("visibility")
    auto_crawl_schedule = request.get("auto_crawl_schedule")
    
    if not new_name and not visibility and auto_crawl_schedule is None:
        return JSONResponse(status_code=400, content={"error": "Name, visibility, or schedule is required"}, media_type="application/json; charset=utf-8")
        
    if visibility and visibility not in ("private", "organization", "public"):
        return JSONResponse(status_code=400, content={"error": "Invalid visibility"}, media_type="application/json; charset=utf-8")

    conn = sqlite3.connect(str(DB_PATH), timeout=30)
    cursor = conn.cursor()
    
    # Try getting the column if it exists to be safe
    try:
        cursor.execute("SELECT id, owner_id, name, visibility, auto_crawl_schedule FROM categories WHERE id = ?", (folder_id,))
    except Exception:
        cursor.execute("SELECT id, owner_id, name, visibility, 'disable' as auto_crawl_schedule FROM categories WHERE id = ?", (folder_id,))

    row = cursor.fetchone()
    if not row:
        conn.close()
        return JSONResponse(status_code=404, content={"error": "Folder not found"}, media_type="application/json; charset=utf-8")
        
    # Any user can modify their own folders; admin can modify all
    if current_user["role"] != "admin" and row[1] != current_user["id"]:
        conn.close()
        return JSONResponse(status_code=403, content={"error": "Not authorized"}, media_type="application/json; charset=utf-8")

    final_name = new_name if new_name else row[2]
    final_vis = visibility if visibility else row[3]
    final_sched = auto_crawl_schedule if auto_crawl_schedule is not None else row[4]

    cursor.execute("UPDATE categories SET name = ?, visibility = ?, auto_crawl_schedule = ? WHERE id = ?", (final_name, final_vis, final_sched, folder_id))

    conn.commit()

    conn.close()



    needs_save = False



    # 1. Update documents using old category name if name changed

    if new_name and new_name != row[2]:

        old_name = row[2]

        for doc in document_status.values():

            if doc.get("category") == old_name:

                doc["category"] = new_name

                needs_save = True



    # 2. Update visibility if it changed

    if visibility and visibility != row[3]:

        for doc in document_status.values():

            if doc.get("category") == final_name:

                doc["visibility"] = final_vis

                needs_save = True



    if needs_save:

        save_status()



    return JSONResponse(content={"id": folder_id, "name": final_name, "visibility": final_vis}, media_type="application/json; charset=utf-8")





@app.delete("/api/folders/{folder_id}")

async def delete_folder(folder_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    """Delete a folder. Documents inside move to parent folder."""

    if folder_id == 1:

        return JSONResponse(status_code=400, content={"error": "Cannot delete root folder"}, media_type="application/json; charset=utf-8")



    conn = sqlite3.connect(str(DB_PATH), timeout=30)

    conn.row_factory = sqlite3.Row

    cursor = conn.cursor()

    cursor.execute("SELECT id, name, parent_id, owner_id FROM categories WHERE id = ?", (folder_id,))

    row = cursor.fetchone()

    if not row:

        conn.close()

        return JSONResponse(status_code=404, content={"error": "Folder not found"}, media_type="application/json; charset=utf-8")

    if current_user["role"] != "admin" and row["owner_id"] != current_user["id"]:

        conn.close()

        return JSONResponse(status_code=403, content={"error": "Not authorized"}, media_type="application/json; charset=utf-8")



    folder_name = row["name"]

    parent_id = row["parent_id"]



    # Find parent category name for document migration

    parent_name = "General"

    if parent_id:

        cursor.execute("SELECT name FROM categories WHERE id = ?", (parent_id,))

        pr = cursor.fetchone()

        if pr:

            parent_name = pr["name"]



    # Move child folders to parent

    cursor.execute("UPDATE categories SET parent_id = ? WHERE parent_id = ?", (parent_id, folder_id))

    # Delete the folder

    cursor.execute("DELETE FROM categories WHERE id = ?", (folder_id,))

    conn.commit()

    conn.close()



    # Move documents from deleted folder to parent folder

    for doc in document_status.values():

        if doc.get("category") == folder_name:

            doc["category"] = parent_name

    save_status()



    return JSONResponse(content={"message": f"Folder '{folder_name}' deleted"}, media_type="application/json; charset=utf-8")



@app.patch("/api/documents/{doc_id}/category")

async def update_document_category(doc_id: str, request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    if doc_id not in document_status:

        return JSONResponse(status_code=404, content={"error": "Document not found"})

        

    doc_info = document_status[doc_id]

    if current_user["role"] != "admin" and doc_info.get("owner_id") != current_user["id"]:

        return JSONResponse(status_code=403, content={"error": "Not authorized to change category of this document"})

        

    new_category = request.get("category", "").strip()

    if not new_category:

        return JSONResponse(status_code=400, content={"error": "Category is required"})

        

    document_status[doc_id]["category"] = new_category

    save_status()

    

    return {"message": "Category updated successfully", "category": new_category}



@app.patch("/api/documents/{doc_id}/active")

async def update_document_active(doc_id: str, request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    if doc_id not in document_status:

        return JSONResponse(status_code=404, content={"error": "Document not found"})

        

    doc_info = document_status[doc_id]

    if current_user["role"] != "admin" and doc_info.get("owner_id") != current_user["id"]:

        return JSONResponse(status_code=403, content={"error": "Not authorized to change active state of this document"})

        

    is_active = request.get("is_active")

    if is_active is None:

        return JSONResponse(status_code=400, content={"error": "is_active is required"})

        

    document_status[doc_id]["is_active"] = bool(is_active)

    save_status()

    

    return {"message": "Document active state updated successfully", "is_active": bool(is_active)}



@app.patch("/api/documents/{doc_id}/visibility")

async def update_document_visibility(doc_id: str, request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    if doc_id not in document_status:

        return JSONResponse(status_code=404, content={"error": "Document not found"})

    

    doc_info = document_status[doc_id]

    

    # Only owner or admin can change visibility

    if current_user["role"] != "admin" and doc_info.get("owner_id") != current_user["id"]:

        return JSONResponse(status_code=403, content={"error": "Not authorized to change visibility of this document"})

    

    new_visibility = request.get("visibility", "").strip()

    if new_visibility not in ("private", "organization", "public") and not new_visibility.startswith("group_"):

        return JSONResponse(status_code=400, content={"error": "visibility must be one of: private, organization, public, or group_<id>"})

    

    document_status[doc_id]["visibility"] = new_visibility

    save_status()

    

    return {"message": "Visibility updated successfully", "visibility": new_visibility}


@app.post("/api/documents/{doc_id}/reindex")
async def reindex_document(doc_id: str, background_tasks: BackgroundTasks, current_user: dict = fastapi.Depends(get_current_user)):
    if doc_id not in document_status:
        return JSONResponse(status_code=404, content={"error": "Document not found"})
        
    doc = document_status[doc_id]
    if doc.get("owner_id") != current_user.get("id") and current_user.get("role") != "admin":
        return JSONResponse(status_code=403, content={"error": "권한이 없습니다."})
        
    from app.rag_engine import cancelled_jobs, index_pdf_async
    cancelled_jobs.discard(doc_id)
    
    safe_name = doc.get("safe_filename", f"{doc_id}_{doc.get('name')}")
    file_path = DOCS_DIR / safe_name
    
    if not os.path.exists(file_path):
        return JSONResponse(status_code=404, content={"error": "원본 파일을 찾을 수 없어 재인덱싱할 수 없습니다."})

    if not str(file_path).lower().endswith(".pdf"):
        return JSONResponse(status_code=400, content={"error": "PDF 파일만 재인덱싱할 수 있습니다. (변환 실패 문서는 삭제 후 재업로드 해주세요)"})

    document_status[doc_id]["status"] = "pending"
    document_status[doc_id]["progress"] = "재인덱싱 대기 중..."
    save_status()
    
    background_tasks.add_task(index_pdf_async, str(file_path), doc_id)
    return {"success": True, "message": "재인덱싱이 시작되었습니다."}

@app.post("/api/documents/{doc_id}/stop")
async def stop_document_indexing(doc_id: str, current_user: dict = fastapi.Depends(get_current_user)):
    if doc_id not in document_status:
        return JSONResponse(status_code=404, content={"error": "Document not found"})
        
    doc = document_status[doc_id]
    if doc.get("owner_id") != current_user.get("id") and current_user.get("role") != "admin":
        return JSONResponse(status_code=403, content={"error": "권한이 없습니다."})
        
    from app.rag_engine import cancelled_jobs
    cancelled_jobs.add(doc_id)
    
    if doc.get("status") in ["pending", "processing"]:
        document_status[doc_id]["progress"] = "중지 명령 수신 (종료 중...)"
        save_status()
        
    return {"success": True, "message": "인덱싱 중지 명령을 전송했습니다."}

@app.delete("/api/documents/{doc_id}")

async def delete_document(doc_id: str, current_user: dict = fastapi.Depends(get_current_user)):

    if doc_id not in document_status:

        return JSONResponse(status_code=404, content={"error": "Document not found"})

        

    doc_info = document_status[doc_id]

    

    # Permission check: Only Admin or Owner can delete

    if current_user["role"] != "admin" and doc_info.get("owner_id") != current_user["id"]:

        return JSONResponse(status_code=403, content={"error": "Not authorized to delete this document"})

        

    safe_filename = doc_info.get("safe_filename")

    

    cancelled_jobs.add(doc_id)

    

    # Remove from memory and file

    del document_status[doc_id]

    save_status()

    

    # Delete PDF file

    if safe_filename:

        pdf_path = DOCS_DIR / safe_filename

        if pdf_path.exists():

            os.remove(pdf_path)

            

    # Delete tree file

    tree_path = TREES_DIR / f"{doc_id}_structure.json"

    if tree_path.exists():

        os.remove(tree_path)

        

    # Delete from SQLite FTS5 Index

    try:

        import sqlite3

        from app.rag_engine import DB_PATH

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("DELETE FROM docs_fts WHERE doc_id = ?", (doc_id,))

        try:
            cursor.execute("DELETE FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
        except Exception:
            pass

        conn.commit()

        conn.close()

    except Exception as e:

        print(f"Error deleting from SQLite index for {doc_id}: {e}")

        

        

    return {"message": "Document deleted successfully"}



class HITLResolveRequest(BaseModel):

    session_id: str

    approved: bool



@app.post("/api/chat/hitl_resolve")

async def hitl_resolve(req: HITLResolveRequest):

    from app.mcp_engine import resolve_hitl_approval

    success = resolve_hitl_approval(req.session_id, req.approved)

    return {"success": success}



@app.post("/api/chat")

async def chat(request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    query = request.get("query")

    active_docs = request.get("active_docs", [])

    history = request.get("history", [])

    agent_id = request.get("agent_id")

    file_paths = request.get("file_paths", [])

    run_sandbox = request.get("run_sandbox", True)

    

    if not query:

        return JSONResponse(status_code=400, content={"error": "Query is required"})

        

    if not agent_id and not active_docs:

        return JSONResponse(status_code=400, content={"error": "At least one document must be active"})

        

    session_id = request.get("session_id")

    if not session_id:

        import uuid

        session_id = str(uuid.uuid4())

        

    try:

        # Save user query to chat_history

        import sqlite3

        from app.rag_engine import DB_PATH

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("INSERT INTO chat_history (session_id, user_id, role, content) VALUES (?, ?, ?, ?)",

                       (session_id, current_user["id"], "user", query))

        conn.commit()

        conn.close()



        # FastAPI StreamingResponse consumes the async generator directly

        async def event_generator():

            async for chunk in chat_inferential_stream(query, active_docs, history, current_user["id"], session_id, agent_id=agent_id, file_paths=file_paths, run_sandbox=run_sandbox):

                if chunk:

                    yield chunk



        return StreamingResponse(

            event_generator(), 

            media_type="text/event-stream",

            headers={

                "Cache-Control": "no-cache",

                "Connection": "keep-alive",

                "X-Accel-Buffering": "no" # Disables buffering in nginx and similar linux proxies

            }

        )

        

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/chat/revise")

async def chat_revise(request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    message_id = request.get("message_id")

    revise_query = request.get("revise_query")

    active_docs = request.get("active_docs", [])

    history = request.get("history", [])

    session_id = request.get("session_id")

    

    if not message_id or not revise_query or not session_id:

        return JSONResponse(status_code=400, content={"error": "message_id, revise_query, and session_id are required"})

        

    try:

        import sqlite3

        import json

        from app.rag_engine import DB_PATH

        from pageindex.utils import ChatGPT_API_async_stream

        

        # 1. Fetch the exact AI response we are revising

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("SELECT content FROM chat_history WHERE id = ? AND session_id = ? AND user_id = ?", 

                       (message_id, session_id, current_user["id"]))

        row = cursor.fetchone()

        

        if not row:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Message not found"})

            

        content = row[0]

        try:

            versions = json.loads(content)

            if isinstance(versions, list) and len(versions) > 0:

                current_assistant_text = versions[-1]["text"]

            else:

                current_assistant_text = content

        except:

            current_assistant_text = content

            

        # 2. Fetch the immediate preceding user query for context

        cursor.execute("SELECT content FROM chat_history WHERE session_id = ? AND user_id = ? AND role = 'user' AND id < ? ORDER BY id DESC LIMIT 1",

                       (session_id, current_user["id"], message_id))

        prev_user_row = cursor.fetchone()

        user_query = prev_user_row[0] if prev_user_row else "?????†ì�Œ"

        conn.close()

        

        agent_id = request.get("agent_id")

        agent_msg = ""

        if agent_id:

            conn = sqlite3.connect(str(DB_PATH), timeout=30)

            cursor = conn.cursor()

            cursor.execute("SELECT name, system_prompt FROM chat_agents WHERE id = ?", (agent_id,))

            agent_row = cursor.fetchone()

            conn.close()

            if agent_row:

                agent_msg = f"?¹ì‹ ?€ ?¬ìš©??ë§žì¶¤???¹ìˆ˜ ?�ì�´?„íŠ¸ '{agent_row[0]}'?…ë‹ˆ?? ê¸°ì¡´ ?µë???ë³´ì™„?????„ëž˜ ?œìŠ¤???„ë¡¬?„íŠ¸??ì§€?œì‚¬??�„ ?¨ê»˜ ê³ ë ¤?˜ì„¸??\n[System Prompt]:\n{agent_row[1]}\n\n"

        

        answer_prompt = f"""{agent_msg}?¹ì‹ ?€ ë¬¸ì„œë¥??„ë¬¸?�ì�´ê³??•í™•?˜ê²Œ ë¶„ì„�?˜ëŠ” AI ?´ì‹œ?¤í„´?¸ìž…?ˆë‹¤.

?¬ìš©?�ê? ?´ì „???�ì„±??AI???µë????€???˜ì • ?�ëŠ” ë³´ì™„???”ì²­?ˆìŠµ?ˆë‹¤. 

?„ëž˜ ?´ìš©??ë°”íƒ•?¼ë¡œ ê¸°ì¡´ ?µë????„ì „???ˆë¡­ê²??‘ì„±??ì£¼ì„¸?? ?¸ë? ì§€?�ì�„ ?¼ìž�??ì§€?´ë‚´ë©??ˆë? ???©ë‹ˆ??



[?¬ìš©?�ì�˜ ?�ëž˜ ì§ˆë¬¸ ë§¥ë�½]

{user_query}



[기존 AI 답변]

{current_assistant_text}



[수정 및 보완 요청]

{revise_query}



CRITICAL INSTRUCTIONS:

1. ?Œê² ?µë‹ˆ???±ì�˜ ?œë¡  ?†ì�´ **?¤ì§� ?˜ì •/ë³´ì™„???ˆë¡œ???µë?ë§Œì�„ ë°”ë¡œ ì¶œë ¥**?˜ì„¸??

2. ì§�ì „ ?µë????¬í•¨?˜ì–´ ?ˆë�˜ ë¬¸ì„œ ì°¸ì¡° ?œê¸° (?? [DocID#Page])ê°€ ? ìš©???•ë³´?¼ë©´ ?¼ì†�?˜ì? ?Šë�„ë¡?? ì???ì£¼ì„¸??

"""

        

        async def event_generator():

            full_reply = ""

            yield json.dumps({"type": "status", "data": "답변 보완 중..."}) + "\\n"

            async for chunk in ChatGPT_API_async_stream(model="gemini-flash-lite-latest", prompt=answer_prompt):

                if chunk:

                    full_reply += chunk

                    yield json.dumps({"type": "chunk", "data": chunk}) + "\\n"

            

            # After streaming is done, append the new version to the existing message_id in the database

            conn = sqlite3.connect(str(DB_PATH), timeout=30)

            cursor = conn.cursor()

            cursor.execute("SELECT content FROM chat_history WHERE id = ? AND session_id = ? AND user_id = ?", 

                           (message_id, session_id, current_user["id"]))

            row = cursor.fetchone()

            

            if row:

                content = row[0]

                try:

                    versions = json.loads(content)

                    if not isinstance(versions, list):

                        versions = [{"query": None, "text": content}]

                except:

                    versions = [{"query": None, "text": content}]

                

                versions.append({"query": revise_query, "text": full_reply})

                new_content = json.dumps(versions, ensure_ascii=False)

                

                cursor.execute("UPDATE chat_history SET content = ? WHERE id = ?", (new_content, message_id))

                conn.commit()

            conn.close()

            # Send the message_id back to frontend to confirm completion for this message

            yield json.dumps({"type": "message_id", "data": message_id}) + "\\n"



        return StreamingResponse(

            event_generator(),

            media_type="text/event-stream",

            headers={

                "Cache-Control": "no-cache",

                "Connection": "keep-alive",

                "X-Accel-Buffering": "no"

            }

        )

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/chat/edit_assistant")

async def edit_assistant_message(request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    """Allow user to manually edit an assistant's message and save it as a new version."""

    message_id = request.get("message_id")

    new_text = request.get("text")

    session_id = request.get("session_id")

    

    if not message_id or not new_text or not session_id:

        return JSONResponse(status_code=400, content={"error": "Missing required fields"})

        

    try:

        import sqlite3

        import json

        from app.rag_engine import DB_PATH

        

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        import tempfile
        import os
        try:
            debug_path = os.path.join(tempfile.gettempdir(), "debug.txt")
            with open(debug_path, "a", encoding="utf-8") as f:
                f.write(f"EDIT CALLED: msg_id='{message_id}', session='{session_id}', user={current_user['id']}\\n")
        except Exception:
            pass



        cursor.execute("SELECT content FROM chat_history WHERE id = ? AND session_id = ? AND user_id = ?", 

                       (message_id, session_id, current_user["id"]))

        row = cursor.fetchone()

        

        if not row:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Message not found"})

            

        content = row[0]

        try:

            versions = json.loads(content)

            if not isinstance(versions, list):

                versions = [{"query": None, "text": content}]

        except:

            versions = [{"query": None, "text": content}]

        

        versions.append({"query": "사용자 직접 수정", "text": new_text})

        new_content = json.dumps(versions, ensure_ascii=False)

        

        cursor.execute("UPDATE chat_history SET content = ? WHERE id = ?", (new_content, message_id))

        conn.commit()

        conn.close()

        

        return {"status": "success"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/export_docx")

async def export_docx(request: dict):

    content = request.get("content")

    if not content:

        return JSONResponse(status_code=400, content={"error": "Content is required"})

        

    try:

        doc = Document()

        doc.add_heading("RAG AI Assistant Response", 0)

        

        def add_markdown_runs(paragraph, text):

            # Parse inline markdown: **bold**, *italic*, and `code`

            tokens = re.split(r'(\*\*.*?\*\*|\*[^*]+\*|`[^`]+`)', text)

            for token in tokens:

                if not token:

                    continue

                if token.startswith('**') and token.endswith('**'):

                    paragraph.add_run(token[2:-2]).bold = True

                elif token.startswith('*') and token.endswith('*'):

                    paragraph.add_run(token[1:-1]).italic = True

                elif token.startswith('`') and token.endswith('`'):

                    r = paragraph.add_run(token[1:-1])

                    # Basic fallback for code styling

                else:

                    paragraph.add_run(token)



        # Better parsing of the markdown content for Word

        lines = content.split('\n')

        for line in lines:

            line = line.replace('\r', '').strip()

            if not line:

                continue

            if line.startswith('### '):

                p = doc.add_heading(level=3)

                add_markdown_runs(p, line[4:].strip())

            elif line.startswith('## '):

                p = doc.add_heading(level=2)

                add_markdown_runs(p, line[3:].strip())

            elif line.startswith('# '):

                p = doc.add_heading(level=1)

                add_markdown_runs(p, line[2:].strip())

            elif line.startswith('- ') or line.startswith('* '):

                p = doc.add_paragraph(style='List Bullet')

                # For `* item` versus `*italic*`, regex handles block * space properly if stripped above

                text_part = line[2:].strip() if line.startswith('- ') or line.startswith('* ') else line

                add_markdown_runs(p, text_part)

            else:

                p = doc.add_paragraph()

                add_markdown_runs(p, line)

        

        stream = io.BytesIO()

        doc.save(stream)

        stream.seek(0)

        

        # Read stream to send as pure bytes response to avoid generator timeout locks

        docx_bytes = stream.getvalue()

        

        headers = {

            'Content-Disposition': 'attachment; filename="rag_response.docx"',

            'Access-Control-Expose-Headers': 'Content-Disposition'

        }

        

        return fastapi.Response(

            content=docx_bytes,

            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",

            headers=headers

        )

        

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": f"Failed to create DOCX: {str(e)}"})



@app.post("/api/export_pptx")

async def export_pptx(request: dict):

    content = request.get("content")

    if not content:

        return JSONResponse(status_code=400, content={"error": "Content is required"})

    try:
        import httpx
        
        API_KEY = "sk-2slides-f310238ab3813cb7e19ca7f80cd7c811c9ae0c892e18d985892860e3948ae0b0"
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }
        
        async with httpx.AsyncClient(timeout=120.0) as client:
            # Generate slides with a premium modern template
            import random
            premium_themes = [
                "st-1763450718138-5utx9lnia", # Black and Gray Gradient 
                "st-1759917935785-nx0z6ae54", # Blue Modern Project
                "st-1761275876731-ojtslyz6e", # Black White Modern Gradient
                "st-1757840073876-sxlvltrs3"  # Green Modern Futuristic
            ]
            theme_id = random.choice(premium_themes)

            enhanced_content = (
                "[PPT 생성 핵심 수칙 - 반드시 준수할 것!]\n"
                "1. 모든 슬라이드의 **제목(Title)은 무조건 10자 이내의 아주 짧은 명사형**으로 작성하세요.\n"
                "   (예: '내부 정보 유출 사고' -> '내부정보 유출', '중국 배후 해킹 그룹 분석' -> '해킹그룹 분석')\n"
                "2. 본문 내용도 최대한 간결한 개조식(Bullet points)으로 작성하여 폰트 짤림 현상(두 줄 겹침)을 방지하세요.\n"
                "3. 한국어 띄어쓰기를 정확히 하고, 문장이 길어져서 레이아웃을 해치지 않도록 강력하게 압축하세요.\n"
                "4. 전반적으로 세련되고 전문적인(Professional) 톤앤매너를 유지하세요.\n\n"
                f"내용:\n{content}"
            )

            payload = {
                "userInput": enhanced_content,
                "themeId": theme_id
            }
            
            gen_res = await client.post("https://2slides.com/api/v1/slides/generate", headers=headers, json=payload)
            if gen_res.status_code != 200:
                raise Exception(f"Failed to generate slides. Status: {gen_res.status_code}, Response: {gen_res.text}")
                
            gen_data = gen_res.json()
            if not gen_data.get("success"):
                raise Exception(f"2Slides generation error: {gen_data.get('error')}")
                
            download_url = gen_data["data"]["downloadUrl"]
            
            file_res = await client.get(download_url)
            if file_res.status_code != 200:
                raise Exception("Failed to download the generated PPTX from 2slides.")
                
            pptx_bytes = file_res.content
            
        resp_headers = {
            'Content-Disposition': 'attachment; filename="AI_Presentation.pptx"',
            'Access-Control-Expose-Headers': 'Content-Disposition'
        }
        
        import fastapi
        return fastapi.Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers=resp_headers
        )
            
    except Exception as e:
        from fastapi.responses import JSONResponse
        return JSONResponse(status_code=500, content={"error": f"Failed to create PPTX via 2Slides: {str(e)}"})


@app.post("/api/export_hwpx")
async def export_hwpx(request: dict):
    content = request.get("content")
    if not content:
        return JSONResponse(status_code=400, content={"error": "Content is required"})
        
    try:
        import tempfile
        import os
        import subprocess
        
        fd, md_path = tempfile.mkstemp(suffix=".md")
        os.close(fd)
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(content)
            
        hwpx_path = md_path.replace(".md", ".hwpx")
        
        result = subprocess.run(
            ["pypandoc-hwpx", md_path, "-o", hwpx_path], 
            capture_output=True, 
            text=True
        )
        
        if result.returncode != 0:
            if os.path.exists(md_path): os.unlink(md_path)
            if os.path.exists(hwpx_path): os.unlink(hwpx_path)
            return JSONResponse(status_code=500, content={"error": f"HWPX Conversion failed: {result.stderr}"})
            
        with open(hwpx_path, "rb") as f:
            hwpx_bytes = f.read()
            
        if os.path.exists(md_path): os.unlink(md_path)
        if os.path.exists(hwpx_path): os.unlink(hwpx_path)
        
        headers = {
            'Content-Disposition': 'attachment; filename="AI_Response.hwpx"',
            'Access-Control-Expose-Headers': 'Content-Disposition'
        }
        
        return fastapi.Response(
            content=hwpx_bytes,
            media_type="application/octet-stream",
            headers=headers
        )
        
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"Failed to create HWPX: {str(e)}"})



# Admin Endpoints for Multi-User

class OrganizationRequest(BaseModel):

    name: str

    parent_id: Optional[int] = None



class UserRequest(BaseModel):

    username: str

    password: str

    organization_id: int

    full_name: Optional[str] = None

    role: str = "user"



@app.post("/api/admin/organizations")

async def create_organization(org: OrganizationRequest, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("INSERT INTO organizations (name, parent_id) VALUES (?, ?)", (org.name, org.parent_id))

        conn.commit()

        org_id = cursor.lastrowid

        conn.close()

        return {"id": org_id, "name": org.name, "parent_id": org.parent_id, "message": "Organization created"}

    except sqlite3.IntegrityError:

        return JSONResponse(status_code=400, content={"error": "Organization already exists"})

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.get("/api/admin/organizations")

async def get_organizations(current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id, name, parent_id FROM organizations")

        rows = cursor.fetchall()

        conn.close()

        orgs = [{"id": r[0], "name": r[1], "parent_id": r[2]} for r in rows]

        return {"organizations": orgs}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



class OrganizationRenameRequest(BaseModel):

    name: str



@app.patch("/api/admin/organizations/{org_id}")

async def rename_organization(org_id: int, org: OrganizationRenameRequest, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("UPDATE organizations SET name = ? WHERE id = ?", (org.name, org_id))

        if cursor.rowcount == 0:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Organization not found"})

        conn.commit()

        conn.close()

        return {"message": "Organization renamed to " + org.name}

    except sqlite3.IntegrityError:

        return JSONResponse(status_code=400, content={"error": "Organization name already exists or invalid data"})

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.delete("/api/admin/organizations/{org_id}")

async def delete_organization(org_id: int, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        # Check if it has children or users before deletion

        cursor.execute("SELECT COUNT(*) FROM users WHERE organization_id = ?", (org_id,))

        if cursor.fetchone()[0] > 0:

            return JSONResponse(status_code=400, content={"error": "Cannot delete organization with assigned users."})

            

        cursor.execute("SELECT COUNT(*) FROM organizations WHERE parent_id = ?", (org_id,))

        if cursor.fetchone()[0] > 0:

            return JSONResponse(status_code=400, content={"error": "Cannot delete organization with child organizations."})

            

        cursor.execute("DELETE FROM organizations WHERE id = ?", (org_id,))

        conn.commit()

        conn.close()

        return {"message": "Organization deleted"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/admin/users")

async def create_user(user: UserRequest, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        from app.auth import get_password_hash

        hashed_pw = get_password_hash(user.password)

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("INSERT INTO users (username, full_name, password_hash, role, organization_id, is_active) VALUES (?, ?, ?, ?, ?, ?)", 

                       (user.username, user.full_name, hashed_pw, user.role, user.organization_id, True))

        conn.commit()

        user_id = cursor.lastrowid

        conn.close()

        return {"id": user_id, "username": user.username, "full_name": user.full_name, "role": user.role, "organization_id": user.organization_id, "message": "User created"}

    except sqlite3.IntegrityError:

        return JSONResponse(status_code=400, content={"error": "Username already exists"})

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.get("/api/admin/users")

async def get_users(current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id, username, full_name, role, organization_id, is_active FROM users")

        rows = cursor.fetchall()

        conn.close()

        users = [{"id": r[0], "username": r[1], "full_name": r[2], "role": r[3], "organization_id": r[4], "is_active": bool(r[5])} for r in rows]

        return {"users": users}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



class UserUpdateRequest(BaseModel):

    is_active: Optional[bool] = None

    role: Optional[str] = None

    organization_id: Optional[int] = None

    password: Optional[str] = None

    full_name: Optional[str] = None



@app.patch("/api/admin/users/{user_id}")

async def update_user(user_id: int, updates: UserUpdateRequest, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        updates_dict = updates.model_dump(exclude_unset=True)

        if not updates_dict:

            return {"message": "No fields to update"}

            

        if "password" in updates_dict:

            from app.auth import get_password_hash

            updates_dict["password_hash"] = get_password_hash(updates_dict.pop("password"))

        

        new_org_id = updates_dict.get("organization_id")

            

        set_clause = ", ".join([f"{k} = ?" for k in updates_dict.keys()])

        values = list(updates_dict.values()) + [user_id]

        

        cursor.execute(f"UPDATE users SET {set_clause} WHERE id = ?", tuple(values))

        conn.commit()

        conn.close()

        

        # If organization changed, update all documents owned by this user

        if new_org_id is not None:

            for doc_info in document_status.values():

                if doc_info.get("owner_id") == user_id:

                    doc_info["organization_id"] = new_org_id

            # Persist the updated status

            save_status()

        

        return {"message": "User updated"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.delete("/api/admin/users/{user_id}")

async def delete_user(user_id: int, current_user: dict = fastapi.Depends(get_current_active_admin)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        # Don't delete the main admin

        cursor.execute("SELECT username FROM users WHERE id = ?", (user_id,))

        row = cursor.fetchone()

        if row and row[0] == "admin":

            return JSONResponse(status_code=400, content={"error": "Cannot delete main admin account"})

            

        cursor.execute("DELETE FROM chat_history WHERE user_id = ?", (user_id,))

        cursor.execute("DELETE FROM users WHERE id = ?", (user_id,))

        conn.commit()

        conn.close()

        return {"message": "User deleted"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.get("/api/admin/settings/llm")
async def get_admin_llm_settings(current_user: dict = fastapi.Depends(get_current_active_admin)):
    from app.rag_engine import get_sys_setting
    settings = {
        "index_llm_model": get_sys_setting("index_llm_model", "gemini-flash-lite-latest"),
        "ocr_llm_model": get_sys_setting("ocr_llm_model", "gemini-flash-lite-latest"),
        "crawl_llm_model": get_sys_setting("crawl_llm_model", "gemini-flash-lite-latest"),
        "chat_vision_llm_model": get_sys_setting("chat_vision_llm_model", "gemini-2.5-flash"),
        "chat_llm_model": get_sys_setting("chat_llm_model", "gemini-flash-lite-latest")
    }
    return JSONResponse(content=settings)

@app.post("/api/admin/settings/llm")
async def post_admin_llm_settings(request: fastapi.Request, current_user: dict = fastapi.Depends(get_current_active_admin)):
    try:
        from app.rag_engine import set_sys_setting
        body = await request.json()
        saved = []
        for key in ["index_llm_model", "ocr_llm_model", "crawl_llm_model", "chat_vision_llm_model", "chat_llm_model"]:
            if key in body:
                if set_sys_setting(key, body[key]):
                    saved.append(key)
        return {"message": "Settings saved successfully", "saved": saved}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.get("/api/admin/documents")

async def admin_get_documents(current_user: dict = fastapi.Depends(get_current_active_admin)):

    docs = []

    for doc_id, info in document_status.items():

        docs.append({

            "id": doc_id,

            "status": info.get("status"),

            "name": info.get("name"),

            "safe_filename": info.get("safe_filename"),

            "visibility": info.get("visibility", "public"),

            "owner_id": info.get("owner_id"),

            "organization_id": info.get("organization_id"),

            "upload_date": info.get("upload_date", "N/A"),

            "page_count": info.get("page_count", "N/A")

        })

    return {"documents": docs}



@app.get("/api/chat/history")

async def get_chat_history(current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT session_id, role, content, timestamp FROM chat_history WHERE user_id = ? ORDER BY timestamp ASC", (current_user["id"],))

        rows = cursor.fetchall()

        conn.close()

        history = [{"session_id": r[0], "role": r[1], "content": r[2], "timestamp": r[3]} for r in rows]

        return {"history": history}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€

# Shared Group Management API

# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€



class SharedGroupMember(BaseModel):

    target_type: str # 'user' or 'organization'

    target_id: int



class SharedGroupCreateUpdate(BaseModel):

    name: str

    members: list[SharedGroupMember]



@app.get("/api/search/members")

async def search_members(q: str, current_user: dict = fastapi.Depends(get_current_user)):

    if not q or len(q.strip()) < 1:

        return {"results": []}

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        results = []

        # Search users

        cursor.execute('''

            SELECT u.id, u.full_name, u.username, o.name 

            FROM users u 

            LEFT JOIN organizations o ON u.organization_id = o.id 

            WHERE u.full_name LIKE ? OR u.username LIKE ? OR o.name LIKE ?

            LIMIT 20

        ''', (f'%{q}%', f'%{q}%', f'%{q}%'))

        user_rows = cursor.fetchall()

        for r in user_rows:

            results.append({

                "target_type": "user",

                "target_id": r[0],

                "name": r[1] or r[2], # full_name or username

                "org_name": r[3] or "?Œì†� ?†ì�Œ",

                "display": f"👤 {r[1] or r[2]} ({r[3] or '?Œì†� ?†ì�Œ'})"

            })

            

        # Search organizations

        cursor.execute('SELECT id, name FROM organizations WHERE name LIKE ? LIMIT 10', (f'%{q}%',))

        org_rows = cursor.fetchall()

        for r in org_rows:

            results.append({

                "target_type": "organization",

                "target_id": r[0],

                "name": r[1],

                "org_name": "ì¡°ì§� ?„ì²´",

                "display": f"?�¢ [{r[1]}] ì¡°ì§� ?„ì²´"

            })

            

        conn.close()

        return {"results": results}

    except Exception as e:

        import traceback

        return JSONResponse(status_code=500, content={"error": str(e) + traceback.format_exc()})



@app.get("/api/shared-groups")

async def get_shared_groups(current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        conn.row_factory = sqlite3.Row

        cursor = conn.cursor()

        

        cursor.execute("SELECT id, name FROM sharing_groups WHERE owner_id = ?", (current_user["id"],))

        groups = cursor.fetchall()

        

        results = []

        for g in groups:

            group_id = g["id"]

            group_name = g["name"]

            

            cursor.execute('''

                SELECT target_type, target_id FROM sharing_group_members WHERE group_id = ?

            ''', (group_id,))

            members = []

            

            for m in cursor.fetchall():

                t_type = m["target_type"]

                t_id = m["target_id"]

                display = ""

                if t_type == "user":

                    c2 = conn.cursor()

                    c2.execute("SELECT u.full_name, u.username, o.name FROM users u LEFT JOIN organizations o ON u.organization_id=o.id WHERE u.id=?", (t_id,))

                    ur = c2.fetchone()

                    if ur:

                        display = f"👤 {ur[0] or ur[1]} ({ur[2] or '?Œì†� ?†ì�Œ'})"

                elif t_type == "organization":

                    c2 = conn.cursor()

                    c2.execute("SELECT name FROM organizations WHERE id=?", (t_id,))

                    mr = c2.fetchone()

                    if mr:

                        display = f"?�¢ [{mr[0]}] ì¡°ì§� ?„ì²´"

                

                if not display:

                    display = f"[?? œ??ê°�ì²´ ID:{t_id}]"

                    

                members.append({

                    "target_type": t_type,

                    "target_id": t_id,

                    "display": display

                })

                

            results.append({

                "id": group_id,

                "name": group_name,

                "members": members

            })

            

        conn.close()

        return {"groups": results}

    except Exception as e:

        import traceback

        return JSONResponse(status_code=500, content={"error": str(e) + traceback.format_exc()})



@app.post("/api/shared-groups")

async def create_shared_group(group: SharedGroupCreateUpdate, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("INSERT INTO sharing_groups (name, owner_id) VALUES (?, ?)", (group.name, current_user["id"]))

        group_id = cursor.lastrowid

        

        for m in group.members:

            cursor.execute("INSERT INTO sharing_group_members (group_id, target_type, target_id) VALUES (?, ?, ?)", 

                           (group_id, m.target_type, m.target_id))

                           

        conn.commit()

        conn.close()

        return {"message": "Group created", "id": group_id}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.put("/api/shared-groups/{group_id}")

async def update_shared_group(group_id: int, group: SharedGroupCreateUpdate, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        # Verify ownership

        cursor.execute("SELECT id FROM sharing_groups WHERE id = ? AND owner_id = ?", (group_id, current_user["id"]))

        if not cursor.fetchone():

            return JSONResponse(status_code=403, content={"error": "Not authorized to edit this group."})

            

        cursor.execute("UPDATE sharing_groups SET name = ? WHERE id = ?", (group.name, group_id))

        

        cursor.execute("DELETE FROM sharing_group_members WHERE group_id = ?", (group_id,))

        for m in group.members:

            cursor.execute("INSERT INTO sharing_group_members (group_id, target_type, target_id) VALUES (?, ?, ?)", 

                           (group_id, m.target_type, m.target_id))

                           

        conn.commit()

        conn.close()

        return {"message": "Group updated"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.delete("/api/shared-groups/{group_id}")

async def delete_shared_group(group_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT id FROM sharing_groups WHERE id = ? AND owner_id = ?", (group_id, current_user["id"]))

        if not cursor.fetchone():

            return JSONResponse(status_code=403, content={"error": "Not authorized to delete this group."})

            

        cursor.execute("DELETE FROM sharing_groups WHERE id = ?", (group_id,))

        cursor.execute("DELETE FROM sharing_group_members WHERE group_id = ?", (group_id,))

        conn.commit()

        conn.close()

        return {"message": "Group deleted"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€

# Chat Agents Management API

# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€



@app.post("/api/agents/generate")

async def generate_agent_config(

    requirements: str = fastapi.Form(...),

    template_file: Optional[fastapi.UploadFile] = fastapi.File(None),

    requires_user_prompt: str = fastapi.Form("false"),

    requires_file_upload: str = fastapi.Form("false"),

    requires_daemon: str = fastapi.Form("false"),

    current_user: dict = fastapi.Depends(get_current_user)

):

    try:

        from pageindex.utils import ChatGPT_API_async

        import json

        import os

        import tempfile

        import shutil

        

        template_context = ""

        

        if template_file and template_file.filename:

            # 1. ?Œë“œë°•ì‹±: ?œí”Œë¦??Œì�¼???„ì‹œ ?´ë�”??ìº�ì‹±?©ë‹ˆ??

            tmp_fd, tmp_path = tempfile.mkstemp(suffix="." + template_file.filename.split('.')[-1])

            os.close(tmp_fd)

            

            try:

                with open(tmp_path, "wb") as f:

                    shutil.copyfileobj(template_file.file, f)

                

                # 2. ë¡œì»¬ ?Œì„œ: DOCX, PPTX êµ¬ì¡°ë¥?ê°•ì œë¡??´ì²´?˜ì—¬ ?¤íŠ¸ë§??¸ë¦¬ë¡?ì¶”ì¶œ?©ë‹ˆ??

                ext = template_file.filename.lower().split('.')[-1]

                extracted_text = ""

                

                if ext in ['docx', 'doc']:

                    import docx

                    doc = docx.Document(tmp_path)

                    texts = [p.text for p in doc.paragraphs if p.text.strip()]

                    # Tables extraction is essential for corporate templates

                    for table in doc.tables:

                        for row in table.rows:

                            for cell in row.cells:

                                texts.extend([p.text for p in cell.paragraphs if p.text.strip()])

                    extracted_text = "\n".join(texts)

                elif ext in ['pptx', 'ppt']:

                    import pptx

                    pres = pptx.Presentation(tmp_path)

                    texts = []

                    for slide in pres.slides:

                        for shape in slide.shapes:

                            if shape.has_text_frame:

                                for p in shape.text_frame.paragraphs:

                                    if p.text.strip():

                                        texts.append(p.text)

                            elif shape.has_table:

                                for row in shape.table.rows:

                                    for cell in row.cells:

                                        if cell.text_frame:

                                            for p in cell.text_frame.paragraphs:

                                                if p.text.strip():

                                                    texts.append(p.text)

                    extracted_text = "\n".join(texts)

                elif ext in ['xlsx', 'csv', 'txt', 'md']:

                    with open(tmp_path, "r", encoding="utf-8", errors="ignore") as f:

                        extracted_text = f.read()

                        

                # 3. LLM ì»¨í…�?¤íŠ¸ ?„ë¡¬?„íŠ¸ ?¸ì �??
                if extracted_text.strip():

                    snippet = extracted_text[:2500]

                    template_context = f"""

===================================================

[사용자 제공 템플릿 문서 구조 미리보기 ({template_file.filename})]

{snippet}

===================================================

* [ë§¤ìš° ì¤‘ìš”] ?¬ìš©?�ê? ?�ë�™ ê²°ê³¼ ?�ì„±???„í•œ ë§ˆìŠ¤???œí”Œë¦??Œì�¼(?? Word, PPT ????ì²¨ë??ˆìŠµ?ˆë‹¤!

* ?¹ì‹ ?€ ë°˜ë“œ?????œí”Œë¦?ë¬¸ì„œ??êµ¬ì¡°ë¥??„ë²½??ë¶„ì„�??'?œí”Œë¦?ê¸°ë°˜ ë¬¸ì„œ ?�ì„± ?�ì�´?„íŠ¸'ë¡??‹íŒ…?˜ì–´???©ë‹ˆ??

* `system_prompt` ?�ëŠ” "?¹ì‹ ?€ ?œê³µ??ë¬¸ì„œ ?œí”Œë¦?{template_file.filename})??Placeholder?¤ì�„ ?¬ìš©?�ì�˜ ?…ë ¥ ?´ìš©?¼ë¡œ ?•ë??˜ê²Œ ì±„ì›Œ?£ì–´ ë¬¸ì„œë¥??„ì„±?˜ëŠ” ?„ë¬¸ê°€?…ë‹ˆ??" ?¼ëŠ” ?´ìš©???µì‹¬ ì§€?œì‚¬??œ¼ë¡?ê°•ë ¥?˜ê²Œ ?¬í•¨?´ì•¼ ?©ë‹ˆ??

* `python_code` 작성 규칙:

  1. ?¸ë? ?Œì�¼(`data.json` ????`open()`?¼ë¡œ ?½ì–´???°ì�´?°ë? ê°€?¸ì˜¤?¤ëŠ” ì½”ë“œë¥??‘ì„±?´ì„œ??**?ˆë?** ???©ë‹ˆ?? (?ŒìŠ¤???˜ê²½?�ì„œ ?„ë¡œê·¸ëž¨??ì¦‰ì‹œ ì¶©ë�Œ?©ë‹ˆ??.

  2. ?�ì�´?„íŠ¸???¬ìš©?�ì????¤ì‹œê°??€??ë§¥ë�½??ê¸°ë°˜?¼ë¡œ ???œí”Œë¦¿ì�˜ ë¹ˆì¹¸??ì±„ì›Œ???©ë‹ˆ?? ?Œì�¼ ?�ë‹¨?�ì„œ `os.environ.get("AGENT_USER_PROMPT", "ê¸°ë³¸ ?ŒìŠ¤???´ìš©...")` ë§¤ê°œë³€?˜ë? ?½ì–´?¤ì„¸??

  3. `from pageindex.utils import ChatGPT_API, extract_json` ???„í�¬???˜ì„¸?? `ChatGPT_API(model="ëª¨ë�¸ëª?, prompt="?„ë¡¬?„íŠ¸") -> str` ?…ë‹ˆ?? ë§¤ê°œë³€??`model`ë¡?`"gemini-2.5-flash"`ë¥??„ë‹¬?˜ê³ , `prompt` ë§¤ê°œë³€?˜ì—�???¤ì§� JSON ?•ì‹�ë§Œì�„ ê°•ì œ?˜ëŠ” ì¶”ì¶œ ?„ë¡¬?„íŠ¸ë¥??„ë‹¬?˜ì„¸??

  4. ?�ìŠ¤???Œì‹±???„í•´ ?ˆë? `json.loads`ë¥??°ì? ë§�ê³ , ë°˜ë“œ??`extracted_dict = extract_json(ê²°ê³¼ë¬¸ìž�??` ???´ìš©???Œì‹±?˜ì„¸?? (````json ë§ˆí�¬?¤ìš´???žì—¬?ˆì–´ ê¸°ë³¸ json ?¼ì�´ë¸ŒëŸ¬ë¦¬ëŠ” ë»—ìŠµ?ˆë‹¤). ?Œì‹± ?¤íŒ¨ ??`{{}}` ê¸°ë³¸ê°’ì�„ ?‹íŒ…?˜ëŠ” ?ˆì™¸ì²˜ë¦¬ë¥?ë°˜ë“œ???¬í•¨?˜ì„¸??

  5. ?œí”Œë¦??Œì�¼???Œë“œ??PPTX??ê²½ìš° **?ˆë?** `open()` ?¼ë¡œ ?�ìŠ¤?¸ì²˜???½ì–´?¤ì�´ì§€ ë§ˆì„¸???¸ì½”???�ëŸ¬ ë°œìƒ�)! ë°˜ë“œ??`sys.argv[1:]` ë°°ì—´???œíšŒ??ì°¾ì? ?œí”Œë¦??Œì�¼(`{ext}`)??`python-docx` ?�ëŠ” `python-pptx` ?¼ì�´ë¸ŒëŸ¬ë¦¬ë¡œ ?ˆì „?˜ê²Œ ë¡œë“œ?˜ì—¬ `extracted_dict` ?ˆì�˜ ?°ì�´?°ë“¤ë¡??´ë? ?�ìŠ¤?¸ë? êµ�ì²´????ê²°ê³¼ ê°�ì²´ë¥?`save()` ?˜ì„¸??

     - `extracted_dict`???¤ì—�??`{{{{}}}}` ê¸°í˜¸ê°€ ?†ìœ¼ë¯€ë¡? ì¹˜í™˜???ŒëŠ” ë°˜ë“œ??`"{{" + key + "}}"` ?•íƒœ(ì¦? ì¤‘ê´„????ê°?ë¡?ë¬¸ì„œ ?´ì�˜ ë¬¸ìž�?´ê³¼ ë§¤ì¹­?œì¼œ ì¹˜í™˜?´ì•¼ ?©ë‹ˆ?? (?ˆë? ì¤‘ê´„????ê°œë§Œ ë§¤ì¹­?˜ì? ë§ˆì„¸??

     - PPTX ì¹˜í™˜ ì£¼ì�˜?¬í•­: ?„í˜• ?ˆì—� ?�ìŠ¤?¸ê? ?ˆê±°????`has_table`)ê°€ ?ˆì�„ ???ˆê³ , ê·¸ë£¹?”ë�œ ?„í˜•(`shape.shape_type == 6`)??ê²½ìš° ?˜ìœ„ ?„í˜•???¬ê??�ìœ¼ë¡??�ìƒ‰?´ì•¼ ë¹ˆì¹¸???“ì¹˜ì§€ ?ŠìŠµ?ˆë‹¤.

     - `{{{{ê¸°í˜¸}}}}` ë¬¸ë²•???¬ëŸ¬ `run` ê°�ì²´ë¡?ìª¼ê°œ?¸ì„œ `replace`ê°€ ??ë¨¹ížˆ??ë²„ê·¸ë¥??ˆë°©?˜ê¸° ?„í•´, ë³µìž¡??run ?¨ìœ„ ?µì œê°€ ?´ë µ?¤ë©´ ?¨ìˆœ??`paragraph.text = paragraph.text.replace(...)` ë¡??µì§¸ë¡???–´?Œì›Œ??ë¬´ë°©?©ë‹ˆ??

  6. ì¹˜í™˜??완료?????¤ë¸Œ?�íŠ¸??ë°˜ë“œ??`ê²°ê³¼_`ê°€ ?‘ë‘�?´ë¡œ ë¶™ì? ?ˆë¡œ???Œì�¼ëª…ìœ¼ë¡??„ìž¬ ?‘ì—… ?”ë ‰? ë¦¬???€?¥í•´???©ë‹ˆ??(?œìŠ¤?œì�´ ?¤ìš´ë¡œë“œ ë§�í�¬ë¥??�ë�™ ?°ë�™?©ë‹ˆ??.

"""

            except Exception as e:

                print(f"[Template Parser Error]: {e}")

            finally:

                try: os.remove(tmp_path)

                except: pass

        prompt_requirement_context = ""

        if requires_user_prompt.lower() == "true" or requires_file_upload.lower() == "true":

            prompt_requirement_context += "\n* [ë§¤ìš° ì¤‘ìš”] `python_code` ?‘ì„± ?? ?¤ì�Œ ?œí”Œë¦?êµ¬ì¡°ë¥?ê¸°ë°˜?¼ë¡œ ?‘ì„±?˜ì„¸??\n```python\nimport os\nimport sys\n"

            if requires_user_prompt.lower() == "true":

                prompt_requirement_context += """

# 절대 sys.argv[1:] 를 사용하여 사용자 질문(검색어, 키워드 등) 파싱하려 하지 마세요. 환경변수를 사용하세요.

user_question = os.environ.get("AGENT_USER_PROMPT", "")

print(f"사용자 질문 수신 완료: {user_question}")

if not user_question.strip():

    print("사용자 질문이 없습니다.")

    sys.exit(1)

"""

            if requires_file_upload.lower() == "true":

                prompt_requirement_context += """

# ì²¨ë??Œì�¼ ê²½ë¡œ ?Œì‹±

if len(sys.argv) > 1:

    file_path = sys.argv[1]

    print(f"ì²¨ë??Œì�¼ ê²½ë¡œ: {file_path}")

else:

    print("ì²¨ë??Œì�¼???…ë¡œ?œë�˜ì§€ ?Šì•˜?µë‹ˆ??")

    sys.exit(1)

"""

            prompt_requirement_context += "```\n"

        

        if requires_daemon.lower() == "true":

            prompt_requirement_context += """

[ë§¤ìš° ì¤‘ìš”: ?°ëª¬(?�ì‹œ ?™ìž‘) ëª¨ë“œ]

???�ì�´?„íŠ¸????ë²??¤í–‰?˜ê³  ?�ë‚˜??ê²ƒì�´ ?„ë‹ˆ?? ?„ë¡œ?¸ìŠ¤ê°€ ì£½ì? ?Šê³  ?�êµ¬?�ìœ¼ë¡?ë°±ê·¸?¼ìš´?œì—�???™ìž‘?´ì•¼ ?©ë‹ˆ??

ë°˜ë“œ??ì½”ë“œ ìµœí•˜?¨ì—� `while True:` ?•íƒœ??ë¬´í•œ ë£¨í”„ë¥??¬í•¨?˜ê³ , ?œìŠ¤???�ì›� ê³ ê°ˆ??ë§‰ê¸° ?„í•´ `import time`ê³?`time.sleep(ì£¼ê¸°)`ë¥?ë£¨í”„ ?ˆì—� ë°˜ë“œ???£ì–´ì£¼ì„¸?? (?? 60ì´ˆë§ˆ???‘ì—… ?¤í–‰)

?„ìš”?˜ë‹¤ë©?`schedule` ëª¨ë“ˆ???¬ìš©?´ì„œ ì£¼ê¸°?�ìœ¼ë¡?ë©”ì�¸ ?¨ìˆ˜ê°€ ?¸ì¶œ?˜ë�„ë¡??‘ì„±?´ë�„ ì¢‹ìŠµ?ˆë‹¤.

"""



        

        prompt = f"""?¹ì‹ ?€ ìµœê³  ?˜ì???AI ?„ë¡¬?„íŠ¸ ?”ì??ˆì–´?´ìž� ?œë‹ˆ???Œì�´???°ì�´??ê³¼í•™?�ìž…?ˆë‹¤.

?¬ìš©?�ì�˜ ?¤ì�Œ ?”êµ¬?¬í•­??ë§žì¶°, ?„ë²½?˜ê²Œ ?™ìž‘?˜ëŠ” AI ?�ì�´?„íŠ¸ ?�ì„±??JSON ?œí”Œë¦¿ì�„ ë§Œë“¤?´ì£¼?¸ìš”.



[사용자 핵심 요구사항]:

{requirements}

{template_context}

{prompt_requirement_context}



[ì¶œë ¥ ê·œì¹™]:

결과는 반드시 아래의 JSON 포맷을 정확히 지켜야 합니다.

{{

    "name": "?�ì�´?„íŠ¸??ì§§ê³  ëª…ë£Œ???´ë¦„ (?? ?‘ì? ?°ì�´??ë¶„ì„�ê°€)",

    "description": "?�ì�´?„íŠ¸????• ???¤ëª…?˜ëŠ” 1~2ë¬¸ìž¥",

    "system_prompt": "AI????• , ?´ì¡°, ì§€?œì‚¬?? ì¶œë ¥ ?•ì‹�??ë§¤ìš° ?�ì„¸?˜ê²Œ ì§€?œí•˜???œìŠ¤???„ë¡¬?„íŠ¸. ì½”ë“œ ?¤í–‰ ê²°ê³¼ê°€ ?ˆë‹¤ë©??´ë? ?´ë–»ê²??´ì„�? ì????¬í•¨.",

    "python_code": "ë§Œì•½ ?Œë�¼ë°”ì�´?¸ê¸‰ ?°ì�´??ì²˜ë¦¬, ?Œì�¼ ë¶„ì„�, ?¸ë? ?œí”Œë¦?ì¹˜í™˜ ???Œì�´??ì½”ë“œ ?¤í–‰???„ìš”?˜ë‹¤ë©??‘ì„±?˜ì„¸?? [ë§¤ìš° ì¤‘ìš”] 'ì²¨ë??Œì�¼'??ì²˜ë¦¬?˜ëŠ” ê²½ìš°???œí•´?œë§Œ `sys.argv[1:]`?�ì„œ ?Œì�¼ ê²½ë¡œë¥??»ìœ¼?¸ìš”! ?¬ìš©?�ì�˜ ?¼ë°˜ ?�ìŠ¤??ì§ˆë¬¸(ê²€?‰ì–´ ???€ ?ˆë? `sys.argv`ë¡?ë°›ì? ë§�ê³  ?˜ê²½ë³€?˜ë? ?¬ìš©?´ì•¼ ?©ë‹ˆ??",

    "requires_file_upload": true ?¹ì? false (?Œì�¼ ?…ë¡œ??ì²˜ë¦¬ê°€ ?„ìš”??ê²½ìš°?�ë§Œ true)

}}

"""

        # Call Gemini with strict JSON validation
        from app.rag_engine import get_sys_setting
        chat_model_cfg = get_sys_setting("chat_llm_model", "gemini-flash-lite-latest")
        raw_json_str = await ChatGPT_API_async(
            model=chat_model_cfg, 

            prompt=prompt, 

            response_mime_type="application/json"

        )

        

        if raw_json_str.startswith("```"):

            raw_json_str = raw_json_str.strip("```json").strip("```")

            

        agent_data = json.loads(raw_json_str.strip())

        return agent_data

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": "Agent generation failed: " + str(e)})



@app.post("/api/agents/modify_code")

async def modify_agent_code(

    original_code: str = fastapi.Form(""),

    modification_prompt: str = fastapi.Form(""),

    requires_user_prompt: str = fastapi.Form("false"),

    requires_file_upload: str = fastapi.Form("false"),

    current_user: dict = fastapi.Depends(get_current_user)

):

    try:

        if not original_code.strip() and not modification_prompt.strip():

            return JSONResponse({"python_code": ""})



        from pageindex.utils import ChatGPT_API_async

        import re



        system_rules = "?¹ì‹ ?€ ìµœê³  ?˜ì????œë‹ˆ???Œì�´??ê°œë°œ?�ìž…?ˆë‹¤. ì£¼ì–´ì§??Œì�´??ì½”ë“œë¥??¬ìš©?�ì�˜ ?˜ì • ?”êµ¬?¬í•­??ë§žê²Œ ë³€ê²½í•˜???„ì „???„ì„±???¨ì�¼ ?Œì�´???¤í�¬ë¦½íŠ¸ë¥??‘ì„±?´ì£¼?¸ìš”. ê¸°ì¡´??? ìš©??ë¹„ì¦ˆ?ˆìŠ¤ ë¡œì§�(ê²€?? ë¶„ì„�, API ?¸ì¶œ ???€ ìµœë???ê·¸ë?ë¡?ë³´ì¡´?´ì•¼ ?©ë‹ˆ??"

        

        rule_prompt = ""

        if requires_user_prompt.lower() == "true":

            rule_prompt += "- ì½”ë“œ ?´ì—�??ë°˜ë“œ???¬ìš©??ì§ˆë¬¸??`os.environ.get('AGENT_USER_PROMPT', '')` ë¡??½ì–´?¤ëŠ” ë¡œì§�???¬í•¨?˜ì–´???©ë‹ˆ?? ?ˆë? `sys.argv[1:]`ë¡?ë¬¸ìž�??ì§ˆë¬¸???Œì‹±?˜ì? ë§ˆì„¸??\n"

        else:

            rule_prompt += "- ë§Œì•½ ê¸°ì¡´ ì½”ë“œ??`AGENT_USER_PROMPT` ?˜ê²½ ë³€?˜ë? ?Œì‹±?˜ëŠ” ë¶€ë¶„ì�´ ?ˆë‹¤ë©??œê±°?˜ì„¸??\n"

            

        if requires_file_upload.lower() == "true":

            rule_prompt += "- ì½”ë“œ ?´ì—�??ë°˜ë“œ??ì²¨ë??Œì�¼ ê²½ë¡œë¥?`sys.argv[1]` ë¡?ë°›ì•„?¤ëŠ” ë¡œì§�???¬í•¨?˜ì–´???©ë‹ˆ??\n"

        else:

            rule_prompt += "- ë§Œì•½ ê¸°ì¡´ ì½”ë“œ??ì²¨ë??Œì�¼ ê²½ë¡œë¥?`sys.argv`ë¡??Œì‹±?˜ëŠ” ë¶€ë¶„ì�´ ?ˆë‹¤ë©??œê±°?˜ì„¸??\n"



        prompt = f"""

{system_rules}



[ê¸°ì¡´ ?Œì�´??ì½”ë“œ]:

```python

{original_code}

```



[수정 요구사항]:

{modification_prompt}



[?„ìˆ˜ ë°˜ì˜� ?�íƒœ ê·œì¹™]:

{rule_prompt}



?¤ì§� ?„ì„±???Œì�´??ì½”ë“œ ë¸”ë¡�ë§?ì¶œë ¥?˜ì„¸?? markdown ?•ì‹�(```python ... ```)???¬ìš©?˜ê³ , ì½”ë“œ ë°–ì�˜ ?´ë– ???¤ëª… ë¬¸êµ¬???�ìŠ¤?¸ë�„ ?ˆë? ì¶œë ¥?˜ì? ë§ˆì„¸??

"""

        

        from app.rag_engine import get_sys_setting
        chat_model_cfg = get_sys_setting("chat_llm_model", "gemini-flash-lite-latest")
        result_text = await ChatGPT_API_async(
            model=chat_model_cfg,

            prompt=prompt

        )

        

        code_result = result_text.strip()

        if "```" in code_result:

            match = re.search(r"```(?:python|py)?(.*?)```", code_result, re.DOTALL | re.IGNORECASE)

            if match:

                code_result = match.group(1).strip()

            else:

                code_result = code_result.replace("```python", "").replace("```py", "").replace("```", "").strip()

            

        return JSONResponse({"python_code": code_result})

    except Exception as e:

        import traceback

        traceback.print_exc()

        return JSONResponse(status_code=500, content={"error": "코드 수정 실패: " + str(e)})



class AgentCreateRequest(BaseModel):

    name: str

    description: Optional[str] = ""

    system_prompt: str

    python_code: Optional[str] = ""

    requires_file_upload: bool = False

    template_filename: Optional[str] = None

    share_scope: str = "PRIVATE"

    agent_type: str = "AUTONOMOUS"

    config: Optional[str] = None

    is_active: bool = False



class AgentUpdateRequest(BaseModel):

    name: Optional[str] = None

    description: Optional[str] = None

    system_prompt: Optional[str] = None

    python_code: Optional[str] = None

    requires_file_upload: Optional[bool] = None

    template_filename: Optional[str] = None

    share_scope: Optional[str] = None

    agent_type: Optional[str] = None

    config: Optional[str] = None

    is_active: Optional[bool] = None



@app.get("/api/agents")

async def get_agents(current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        # visible if: owner OR ALL OR (ORG and matching org)

        org_id = current_user.get("organization_id")

        

        # Ensure column exists

        try:

            cursor.execute("ALTER TABLE chat_agents ADD COLUMN template_filename TEXT")

        except:

            pass

            

        user_id = current_user["id"]
        org_id = current_user.get("organization_id")
        
        grp_query = """
        SELECT id FROM sharing_groups WHERE owner_id = ?
        UNION
        SELECT group_id FROM sharing_group_members WHERE target_type = 'user' AND target_id = ?
        """
        grp_params = [user_id, user_id]
        if org_id is not None:
             grp_query += " UNION SELECT group_id FROM sharing_group_members WHERE target_type = 'organization' AND target_id = ?"
             grp_params.append(org_id)
             
        cursor.execute(grp_query, tuple(grp_params))
        user_groups = [r[0] for r in cursor.fetchall()]
        
        query = """
            SELECT id, name, description, system_prompt, python_code, requires_file_upload, user_id, organization_id, share_scope, created_at, updated_at, template_filename, agent_type, config, is_active
            FROM chat_agents
            WHERE user_id = ?
               OR share_scope = 'ALL'
        """
        params = [user_id]

        if org_id is not None:
            query += " OR (share_scope = 'ORG' AND organization_id = ?)"
            params.append(org_id)

        if user_groups:
            query += " OR share_scope IN (" + ",".join(["?"]*len(user_groups)) + ")"
            params.extend([f"group_{g}" for g in user_groups])
        cursor.execute(query, tuple(params))

        rows = cursor.fetchall()

        conn.close()

        

        agents = []

        for r in rows:

            agents.append({

                "id": r[0],

                "name": r[1],

                "description": r[2],

                "system_prompt": r[3],

                "python_code": r[4],

                "requires_file_upload": bool(r[5]),

                "user_id": r[6],

                "organization_id": r[7],

                "share_scope": r[8],

                "created_at": r[9],

                "updated_at": r[10],

                "template_filename": r[11],

                "agent_type": r[12],

                "config": r[13],

                "is_active": bool(r[14])

            })

        return {"agents": agents}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/agents")

async def create_agent(agent: AgentCreateRequest, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        if hasattr(agent, "share_scope"):
            print(f"!!! DEBUG_CREATE_SHARE_SCOPE: type={type(agent.share_scope)}, repr={repr(agent.share_scope)} !!!", flush=True)

            

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("""

            INSERT INTO chat_agents 

            (name, description, system_prompt, python_code, requires_file_upload, template_filename, user_id, organization_id, share_scope, agent_type, config, is_active)

            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)

        """, (

            agent.name, 

            agent.description, 

            agent.system_prompt, 

            agent.python_code, 

            agent.requires_file_upload, 

            agent.template_filename,

            current_user["id"], 

            current_user.get("organization_id"), 

            agent.share_scope,

            "AUTONOMOUS",

            agent.config,

            agent.is_active

        ))

        

        conn.commit()

        agent_id = cursor.lastrowid

        conn.close()

        



        return {"id": agent_id, "message": "Agent created successfully"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.put("/api/agents/{agent_id}")

async def update_agent(agent_id: int, updates: AgentUpdateRequest, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("SELECT user_id FROM chat_agents WHERE id = ?", (agent_id,))

        row = cursor.fetchone()

        if not row:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Agent not found"})

            

        # Admin or Owner

        if row[0] != current_user["id"] and current_user.get("role") != "admin":

            conn.close()

            return JSONResponse(status_code=403, content={"error": "Not authorized to edit this agent"})

            

        updates_dict = updates.model_dump(exclude_unset=True)

        if not updates_dict:

            conn.close()

            return {"message": "No fields to update"}

            

        if "share_scope" in updates_dict:
            print(f"!!! DEBUG_INCOMING_SHARE_SCOPE: type={type(updates_dict['share_scope'])}, repr={repr(updates_dict['share_scope'])} !!!", flush=True)

            

        updates_dict["agent_type"] = "AUTONOMOUS"

        updates_dict["updated_at"] = "CURRENT_TIMESTAMP" 

        

        set_clause = ", ".join([f"{k} = ?" if k != "updated_at" else f"{k} = CURRENT_TIMESTAMP" for k in updates_dict.keys()])

        values = [v for k, v in updates_dict.items() if k != "updated_at"] + [agent_id]

        

        cursor.execute(f"UPDATE chat_agents SET {set_clause} WHERE id = ?", tuple(values))

        conn.commit()

        

        cursor.execute("SELECT agent_type, is_active, config FROM chat_agents WHERE id = ?", (agent_id,))

        updated_row = cursor.fetchone()

        conn.close()

        



        return {"message": "Agent updated successfully"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.delete("/api/agents/{agent_id}")

async def delete_agent(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("SELECT user_id FROM chat_agents WHERE id = ?", (agent_id,))

        row = cursor.fetchone()

        if not row:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Agent not found"})

            

        if row[0] != current_user["id"] and current_user.get("role") != "admin":

            conn.close()

            return JSONResponse(status_code=403, content={"error": "Not authorized to delete this agent"})

            

        cursor.execute("DELETE FROM chat_agents WHERE id = ?", (agent_id,))

        conn.commit()

        conn.close()

        

        asyncio.create_task(mcp_manager.stop_mcp_server(agent_id))

        from app.daemon_engine import daemon_manager

        daemon_manager.stop_daemon(agent_id)

        

        try:

            import os, shutil

            from app.sandbox import VENVS_DIR

            venv_path = os.path.join(VENVS_DIR, str(agent_id))

            if os.path.exists(venv_path):

                shutil.rmtree(venv_path, ignore_errors=True)

        except Exception as ve:

            print(f"Failed to clear venv on agent delete: {ve}")

        

        return {"message": "Agent deleted successfully"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/agents/{agent_id}/sandbox/reset")

async def reset_agent_sandbox(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("SELECT user_id FROM chat_agents WHERE id = ?", (agent_id,))

        row = cursor.fetchone()

        conn.close()

        

        if not row:

            return JSONResponse(status_code=404, content={"error": "Agent not found"})

            

        if row[0] != current_user["id"] and current_user.get("role") != "admin":

            return JSONResponse(status_code=403, content={"error": "Not authorized to reset this sandbox"})

            

        import os, shutil

        from app.sandbox import VENVS_DIR

        venv_path = os.path.join(VENVS_DIR, str(agent_id))

        

        if os.path.exists(venv_path):

            shutil.rmtree(venv_path, ignore_errors=True)

            return {"message": "Sandbox environment data cleared.", "reset": True}

        else:

            return {"message": "Sandbox environment is already empty.", "reset": False}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



from pydantic import BaseModel
class SandboxManualCommand(BaseModel):
    command: str

@app.post("/api/agents/{agent_id}/sandbox/terminal")
async def run_sandbox_terminal(agent_id: int, request: SandboxManualCommand, current_user: dict = fastapi.Depends(get_current_user)):
    try:
        from app.sandbox import execute_terminal_command
        import sqlite3
        from app.rag_engine import DB_PATH
        import os
        
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("SELECT id, user_id FROM chat_agents WHERE id=?", (agent_id,))
        agent = cursor.fetchone()
        conn.close()
        
        if not agent:
            return fastapi.responses.JSONResponse(status_code=404, content={"error": "Agent not found."})
            
        owner_id = agent[1]
        
        if owner_id != current_user["id"] and current_user.get("role") != "admin":
            return fastapi.responses.JSONResponse(status_code=403, content={"error": "Unauthorized to modify this agent's sandbox."})
            
        result = execute_terminal_command(agent_id, request.command)
        return result
        
    except Exception as e:
        return fastapi.responses.JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/agents/{agent_id}/sandbox/cancel")
async def cancel_sandbox_execution_route(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):
    try:
        from app.sandbox import cancel_sandbox_execution
        import sqlite3
        from app.rag_engine import DB_PATH
        
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("SELECT id, user_id FROM chat_agents WHERE id=?", (agent_id,))
        agent = cursor.fetchone()
        conn.close()
        
        if not agent:
            return fastapi.responses.JSONResponse(status_code=404, content={"error": "Agent not found."})
            
        owner_id = agent[1]
        
        if owner_id != current_user["id"] and current_user.get("role") != "admin":
            return fastapi.responses.JSONResponse(status_code=403, content={"error": "Unauthorized."})
            
        success = cancel_sandbox_execution(agent_id)
        if success:
            return {"message": "Execution cancelled successfully."}
        else:
            return {"message": "No active execution found or could not be cancelled."}
            
    except Exception as e:
        return fastapi.responses.JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/agents/test")
async def test_agent_code(
    # fastapi Form overrides default JSON when processing File uploads simultaneously
    python_code: str = fastapi.Form(default=""),
    template_file: Optional[fastapi.UploadFile] = fastapi.File(default=None),
    test_files: List[fastapi.UploadFile] = fastapi.File(default=[]),
    test_args: Optional[str] = fastapi.Form(default=None),
    agent_id: Optional[int] = fastapi.Form(default=None),
    current_user: dict = fastapi.Depends(get_current_user)
):
    try:
        import os
        import tempfile
        import shutil
        import asyncio
        import shlex
        from app.sandbox import execute_agent_code
        
        local_test_files = []
        temp_file_paths = []
        
        if test_files:
            for tf in test_files:
                if tf and tf.filename:
                    tmpdir = tempfile.mkdtemp()
                    dest_path = os.path.join(tmpdir, tf.filename)
                    with open(dest_path, "wb") as f:
                        shutil.copyfileobj(tf.file, f)
                    local_test_files.append(dest_path)
                    temp_file_paths.append(tmpdir)
                    
        if template_file and template_file.filename:
            tmpdir = tempfile.mkdtemp()
            dest_path = os.path.join(tmpdir, template_file.filename)
            with open(dest_path, "wb") as f:
                shutil.copyfileobj(template_file.file, f)
            local_test_files.append(dest_path)
            temp_file_paths.append(tmpdir)
            
        custom_args = []
        if test_args and test_args.strip():
            custom_args = shlex.split(test_args.strip())
            
        result = await asyncio.to_thread(execute_agent_code, python_code, local_test_files, "", agent_id, custom_args)
        
        # Cleanup injected temp test files after sandbox run completes

        for d in temp_file_paths:

            try:

                shutil.rmtree(d)

            except:

                pass

                

        return result

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/agents/{agent_id}/daemon/start")

async def start_agent_daemon(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute("SELECT user_id, python_code FROM chat_agents WHERE id = ?", (agent_id,))

        row = cursor.fetchone()

        conn.close()

        if not row:

            return JSONResponse(status_code=404, content={"error": "Agent not found"})

            

        user_id, python_code = row

        if user_id != current_user["id"] and current_user.get("role") != "admin":

            return JSONResponse(status_code=403, content={"error": "Not authorized to start this agent"})

            

        from app.daemon_engine import daemon_manager

        success = daemon_manager.start_daemon(agent_id, python_code)

        if success:
            try:
                import json
                conn_db = sqlite3.connect(str(DB_PATH), timeout=30)
                c_cursor = conn_db.cursor()
                c_cursor.execute("SELECT config FROM chat_agents WHERE id = ?", (agent_id,))
                cfg_txt_row = c_cursor.fetchone()
                if cfg_txt_row and cfg_txt_row[0]:
                    cfg = json.loads(cfg_txt_row[0])
                    if not cfg.get("cfg-auto-daemon"):
                        cfg["cfg-auto-daemon"] = True
                        c_cursor.execute("UPDATE chat_agents SET config = ? WHERE id = ?", (json.dumps(cfg, ensure_ascii=False), agent_id))
                        conn_db.commit()
                conn_db.close()
            except Exception: pass

            return {"message": "Daemon started successfully"}

        else:

            return JSONResponse(status_code=500, content={"error": "Failed to start daemon process"})

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/agents/{agent_id}/daemon/stop")

async def stop_agent_daemon(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        from app.daemon_engine import daemon_manager

        success = daemon_manager.stop_daemon(agent_id)
        
        try:
            import sqlite3, json
            from app.rag_engine import DB_PATH
            conn_db = sqlite3.connect(str(DB_PATH), timeout=30)
            c_cursor = conn_db.cursor()
            c_cursor.execute("SELECT config FROM chat_agents WHERE id = ?", (agent_id,))
            cfg_txt_row = c_cursor.fetchone()
            if cfg_txt_row and cfg_txt_row[0]:
                cfg = json.loads(cfg_txt_row[0])
                if cfg.get("cfg-auto-daemon"):
                    cfg["cfg-auto-daemon"] = False
                    c_cursor.execute("UPDATE chat_agents SET config = ? WHERE id = ?", (json.dumps(cfg, ensure_ascii=False), agent_id))
                    conn_db.commit()
            conn_db.close()
        except Exception: pass

        if success:

            return {"message": "Daemon stopped successfully"}

        else:

             return {"message": "Daemon was not running or failed to stop"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.get("/api/agents/{agent_id}/daemon/status")

async def get_agent_daemon_status(agent_id: int, current_user: dict = fastapi.Depends(get_current_user)):

    try:

        from app.daemon_engine import daemon_manager

        status = daemon_manager.get_status(agent_id)

        logs = daemon_manager.get_logs(agent_id)

        status["logs"] = logs

        return status

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/agents/{agent_id}/template")

async def upload_agent_template(

    agent_id: int, 

    template_file: fastapi.UploadFile = fastapi.File(...), 

    current_user: dict = fastapi.Depends(get_current_user)

):

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        

        cursor.execute("SELECT user_id FROM chat_agents WHERE id = ?", (agent_id,))

        row = cursor.fetchone()

        if not row:

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Agent not found"})

            

        if row[0] != current_user["id"] and current_user.get("role") != "admin":

            conn.close()

            return JSONResponse(status_code=403, content={"error": "Not authorized to upload template to this agent"})

            

        import uuid

        import os

        from app.sandbox import TEMPLATES_DIR

        

        # Ensure secure unique filename mapping

        filename = f"{agent_id}_{uuid.uuid4().hex[:8]}_{template_file.filename}"

        dest_path = os.path.join(TEMPLATES_DIR, filename)

        

        with open(dest_path, "wb") as f:

            import shutil

            shutil.copyfileobj(template_file.file, f)

            

        cursor.execute("UPDATE chat_agents SET template_filename = ? WHERE id = ?", (filename, agent_id))

        conn.commit()

        conn.close()

        

        return {"message": "Template uploaded successfully", "template_filename": filename}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€

# Chat Session Management API

# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€



class ChatSessionPinRequest(BaseModel):
    is_pinned: bool

@app.patch("/api/chat/sessions/{session_id}/pin")
async def toggle_chat_session_pin(session_id: str, request: ChatSessionPinRequest, current_user: dict = fastapi.Depends(get_current_user)):
    try:
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute(
            "UPDATE chat_sessions SET is_pinned = ? WHERE id = ? AND user_id = ?",
            (request.is_pinned, session_id, current_user["id"])
        )
        conn.commit()
        conn.close()
        return {"status": "success", "is_pinned": request.is_pinned}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.get("/api/chat/sessions")

async def get_chat_sessions(current_user: dict = fastapi.Depends(get_current_user)):

    """Return all chat sessions for the current user, newest first."""

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute(

            """SELECT id, title, created_at, updated_at, COALESCE(is_pinned, 0) as is_pinned
               FROM chat_sessions 
               WHERE user_id = ? 
                 AND EXISTS (SELECT 1 FROM chat_history WHERE chat_history.session_id = chat_sessions.id)
               ORDER BY is_pinned DESC, updated_at DESC""",

            (current_user["id"],)

        )

        rows = cursor.fetchall()

        conn.close()

        sessions = [{"id": r[0], "title": r[1], "created_at": r[2], "updated_at": r[3], "is_pinned": bool(r[4])} for r in rows]

        return {"sessions": sessions}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})





@app.post("/api/chat/sessions")

async def create_chat_session(current_user: dict = fastapi.Depends(get_current_user)):

    """Create a new chat session and return its id."""

    try:

        import uuid as _uuid

        session_id = str(_uuid.uuid4())

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute(

            "INSERT INTO chat_sessions (id, user_id, title) VALUES (?, ?, ?)",

            (session_id, current_user["id"], "???")

        )

        conn.commit()

        conn.close()

        return {"session_id": session_id}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})





@app.delete("/api/chat/sessions/{session_id}")

async def delete_chat_session(session_id: str, current_user: dict = fastapi.Depends(get_current_user)):

    """Delete a session and all its messages."""

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        # Verify ownership

        cursor.execute("SELECT id FROM chat_sessions WHERE id = ? AND user_id = ?", (session_id, current_user["id"]))

        if not cursor.fetchone():

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Session not found"})

        cursor.execute("DELETE FROM chat_history WHERE session_id = ? AND user_id = ?", (session_id, current_user["id"]))

        cursor.execute("DELETE FROM chat_sessions WHERE id = ? AND user_id = ?", (session_id, current_user["id"]))

        conn.commit()

        conn.close()

        return {"message": "Session deleted"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})





@app.get("/api/chat/sessions/{session_id}/messages")

async def get_session_messages(session_id: str, current_user: dict = fastapi.Depends(get_current_user)):

    """Return all messages for a specific session."""

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute(

            "SELECT id, role, content, timestamp FROM chat_history WHERE session_id = ? AND user_id = ? ORDER BY timestamp ASC",

            (session_id, current_user["id"])

        )

        rows = cursor.fetchall()

        conn.close()

        messages = [{"id": r[0], "role": r[1], "content": r[2], "timestamp": r[3]} for r in rows]

        return {"messages": messages}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})





@app.patch("/api/chat/sessions/{session_id}/title")

async def update_session_title(session_id: str, current_user: dict = fastapi.Depends(get_current_user)):

    """Auto-generate session title from its first user message."""

    try:

        conn = sqlite3.connect(str(DB_PATH), timeout=30)

        cursor = conn.cursor()

        cursor.execute(

            "SELECT id FROM chat_sessions WHERE id = ? AND user_id = ?",

            (session_id, current_user["id"])

        )

        if not cursor.fetchone():

            conn.close()

            return JSONResponse(status_code=404, content={"error": "Session not found"})

        # Get first user message as title base

        cursor.execute(

            "SELECT content FROM chat_history WHERE session_id = ? AND user_id = ? AND role = 'user' ORDER BY timestamp ASC LIMIT 1",

            (session_id, current_user["id"])

        )

        row = cursor.fetchone()

        if row:

            title = row[0][:40].replace('\n', ' ')

            cursor.execute(

                "UPDATE chat_sessions SET title = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",

                (title, session_id)

            )

            conn.commit()

        conn.close()

        return {"title": title if row else "새 대화"}

    except Exception as e:

        return JSONResponse(status_code=500, content={"error": str(e)})



# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€

# Branding API

# ?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€?€



@app.get("/api/branding")

async def get_branding():

    """Public endpoint ??anyone (even unauthenticated) can fetch branding info for the login screen."""

    return load_branding()





class BrandingRequest(BaseModel):
    company_name: str = "추론형 RAG"
    tagline: Optional[str] = ""
    logo_type: Optional[str] = "square"





@app.post("/api/admin/branding")

async def update_branding(req: BrandingRequest, current_user: dict = fastapi.Depends(get_current_active_admin)):

    branding = load_branding()

    branding["company_name"] = req.company_name.strip() or "추론형 RAG"
    branding["tagline"] = req.tagline.strip() if req.tagline else ""
    branding["logo_type"] = req.logo_type if req.logo_type in ["square", "rectangle"] else "square"

    save_branding(branding)

    return {"message": "ë¸Œëžœ???�ìŠ¤?¸ê? ?€?¥ë�˜?ˆìŠµ?ˆë‹¤.", "branding": branding}





@app.post("/api/admin/branding/logo")

async def upload_branding_logo(file: UploadFile = File(...), current_user: dict = fastapi.Depends(get_current_active_admin)):

    # Validate file type

    allowed_types = {"image/png", "image/jpeg", "image/jpg", "image/svg+xml", "image/webp", "image/gif"}

    if file.content_type not in allowed_types:

        return JSONResponse(status_code=400, content={"error": "PNG, JPG, SVG, WebP, GIF 이미지만 업로드 가능합니다."})

    

    # Read and enforce 2 MB size limit

    content = await file.read()

    if len(content) > 2 * 1024 * 1024:

        return JSONResponse(status_code=400, content={"error": "?´ë?ì§€ ?Œì¼ ?¬ê¸°??2MBë¥?ì´ˆê³¼?????†ìŠµ?ˆë‹¤."})

    

    # Save with a fixed filename so it's easy to serve (overwrite each time)

    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else "png"

    logo_filename = f"logo.{ext}"

    logo_path = LOGOS_DIR / logo_filename

    

    # Remove old logos first

    for old in LOGOS_DIR.iterdir():

        if old.is_file():

            old.unlink()

    

    with open(logo_path, "wb") as f:

        f.write(content)

    

    branding = load_branding()

    branding["logo_url"] = f"/logos/{logo_filename}"

    save_branding(branding)

    return {"message": "ë¡œê³ ê°€ ?…ë¡œ?œë˜?ˆìŠµ?ˆë‹¤.", "logo_url": branding["logo_url"]}





@app.delete("/api/admin/branding/logo")

async def delete_branding_logo(current_user: dict = fastapi.Depends(get_current_active_admin)):

    branding = load_branding()

    branding["logo_url"] = ""

    save_branding(branding)

    for old in LOGOS_DIR.iterdir():

        if old.is_file():

            old.unlink()

    return {"message": "로고가 삭제되었습니다.", "branding": branding}





@app.post("/api/search/documents")

async def search_documents_endpoint(request: dict, current_user: dict = fastapi.Depends(get_current_user)):

    query = request.get("query", "").strip()

    active_docs = request.get("active_docs", [])

    

    if not query:

        return JSONResponse(status_code=400, content={"error": "Query is required"})

        

    if not active_docs:

        return JSONResponse(status_code=400, content={"error": "No documents selected for search"})

        

    try:

        from app.rag_engine import search_documents_semantic

        docs_raw = await search_documents_semantic(query, active_docs)

        

        seen_docs = set()

        docs_out = []

        for d in docs_raw:

            if d["id"] not in seen_docs:

                seen_docs.add(d["id"])

                docs_out.append(d)

                

        return {"documents": docs_out}

        

    except Exception as e:

        import traceback

        return JSONResponse(status_code=500, content={"error": str(e), "trace": traceback.format_exc()})




class CrawlRequest(BaseModel):
    url: str
    site_name: str
    folder_id: Optional[int] = None
    max_depth: Optional[int] = 1
    max_pages: Optional[int] = 50
    crawl_type: Optional[str] = "spa"
    strategy: Optional[str] = "bfs"
    restrict_path: Optional[bool] = False
    login_id: Optional[str] = None
    login_pw: Optional[str] = None
    search_keyword: Optional[str] = None
    use_ai_extraction: Optional[bool] = False
    ai_extraction_prompt: Optional[str] = None
    clear_existing: Optional[bool] = True

def background_crawl_website(doc_id: str, req: CrawlRequest, owner_id: int, org_id: int, category_name: str, visibility: str):
    from app.crawler import crawl_spa_sync_wrapper
    from app.rag_engine import document_status, save_status, DB_PATH, cancelled_jobs, TREES_DIR
    import sqlite3
    import uuid
    import time
    from pageindex.utils import ChatGPT_API_async
    import asyncio
    import json

    try:
        # 진행 중으로 상태 업데이트
        if doc_id not in document_status: return
        document_status[doc_id]["progress"] = "페이지 탐색 및 수집 중..."
        document_status[doc_id]["progress_percent"] = 20
        save_status()

        is_clear_existing = getattr(req, 'clear_existing', True)
        existing_cache = {}
        
        try:
            conn_cache = sqlite3.connect(str(DB_PATH), timeout=30)
            cur_cache = conn_cache.cursor()
            if not is_clear_existing:
                cur_cache.execute("SELECT url, content_hash FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
                for row in cur_cache.fetchall():
                    existing_cache[row[0]] = row[1]
            conn_cache.close()
        except Exception as e:
            print(f"Error reading cache: {e}")

        # 크롤링 수행
        # 참고: crawl_spa_sync_wrapper 쪽에 is_cancelled 콜백을 붙일 수도 있지만,
        # Playwright를 쓰는 특성상 한방에 끝나기를 기다리게 됩니다.
        results = crawl_spa_sync_wrapper(
            req.url, req.max_depth, req.max_pages, doc_id, req.login_id, req.login_pw, req.search_keyword, req.crawl_type, req.strategy, req.restrict_path, req.use_ai_extraction, req.ai_extraction_prompt, existing_cache
        )
        
        # 중간 취소 여부 확인
        is_cancelled = doc_id in cancelled_jobs

        if not results:
            conn = sqlite3.connect(str(DB_PATH), timeout=30)
            cursor = conn.cursor()
            cursor.execute("SELECT COUNT(DISTINCT url) FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
            total_pages = cursor.fetchone()[0]
            conn.close()

            document_status[doc_id]["status"] = "ready"
            document_status[doc_id]["progress"] = "변경된 문서 없음 (수집 완료)"
            document_status[doc_id]["progress_percent"] = 100
            document_status[doc_id]["page_count"] = total_pages
            document_status[doc_id]["updated_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            save_status()
            return
            
        document_status[doc_id]["progress"] = "부분 수집됨 - 청크 분할 및 인덱싱 중..." if is_cancelled else "청크 분할 및 인덱싱 중..."
        document_status[doc_id]["progress_percent"] = 70
        save_status()
        # 텍스트 스플리터
        def split_text(text: str, chunk_size: int = 2000, chunk_overlap: int = 300) -> list[str]:
            chunks = []
            if not text: return chunks
            length = len(text)
            for i in range(0, length, chunk_size - chunk_overlap):
                chunk = text[i:i + chunk_size]
                if chunk: chunks.append(chunk)
                if i + chunk_size >= length: break
            return chunks

        import hashlib
        
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        
        is_clear_existing = getattr(req, 'clear_existing', True)
        
        tree_path = TREES_DIR / f"{doc_id}_structure.json"
        
        doc_name = req.site_name
        current_structure = []
        if not is_clear_existing and tree_path.exists():
            try:
                with open(tree_path, "r", encoding="utf-8") as f:
                    old_tree = json.load(f)
                    current_structure = old_tree.get("structure", [])
                    doc_name = old_tree.get("doc_name", doc_name)
            except Exception:
                pass
        
        # 기존 옵션에 따라 전체 초기화
        if is_clear_existing:
            cursor.execute("DELETE FROM docs_fts WHERE doc_id = ?", (doc_id,))
            cursor.execute("DELETE FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
            current_structure = []
        cursor.execute("SELECT url, content_hash FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
        cache_map = {row[0]: row[1] for row in cursor.fetchall()}
        conn.commit()
        conn.close()
        
        chunk_count = 0
        current_urls = set()
        pages_to_summarize = []
        
        for page in results:
            content = page.get("content", "")
            page_url = page.get("url", "")
            if not page_url: continue
            
            current_urls.add(page_url)
            content_hash = hashlib.md5(content.encode('utf-8', errors='ignore')).hexdigest()
            
            if cache_map.get(page_url) == content_hash and not is_clear_existing:
                continue

            pages_to_summarize.append(page)

        async def summarize_and_index_all(pages):
            db_lock = asyncio.Lock()
            # sLLM 엔진의 동시 접속 한도가 5개이므로 크롤링 백그라운드 작업은 2개로 제한하여 병목 방지
            sem = asyncio.Semaphore(2)
            
            async def process_page(pg):
                if doc_id in cancelled_jobs: return None
                async with sem:
                    if doc_id in cancelled_jobs: return None
                    page_url = pg.get("url", "")
                    content = pg.get("content", "")
                
                    # 1. LLM Summary
                    prompt = f"다음 웹 페이지의 본문(텍스트) 핵심 내용만 2~3문장 이내로 명확하게 요약하세요. 사이드바 내용이나 메뉴 링크 등 사이트 전반의 불필요한 노이즈는 제외하고, 본문의 고유한 내용에 집중하세요:\n\n{content[:5000]}"
                    try:
                        from app.rag_engine import get_sys_setting
                        crawl_model_cfg = get_sys_setting("crawl_llm_model", "gemini-flash-lite-latest")
                        res = await ChatGPT_API_async(model=crawl_model_cfg, prompt=prompt)
                        pg['summary'] = res.strip()
                    except Exception as e:
                        print(f"[{doc_id}] LLM Summary failed for {page_url}: {e}")
                        pg['summary'] = "요약을 생성하지 못했습니다."
                        
                    # LLM 부하를 조절하기 위해 요청 간 약간의 지연 추가
                    await asyncio.sleep(0.5)
                    
                    # 2. Instant DB Indexing
                    content_hash = hashlib.md5(content.encode('utf-8', errors='ignore')).hexdigest()
                    chunks = split_text(content, chunk_size=3000, chunk_overlap=300)
                    node_ids_for_page = []
                    
                    async with db_lock:
                        conn_local = sqlite3.connect(str(DB_PATH), timeout=10)
                        cur_local = conn_local.cursor()
                        
                        if not is_clear_existing:
                            cur_local.execute("DELETE FROM docs_fts WHERE doc_id = ? AND page_num = ?", (doc_id, page_url))
                            
                        for chunk in chunks:
                            node_id = str(uuid.uuid4())
                            node_ids_for_page.append(node_id)
                            cur_local.execute('''INSERT INTO docs_fts (doc_id, node_id, title, text_content, page_num)
                                VALUES (?, ?, ?, ?, ?)''',
                                (doc_id, node_id, pg['title'] or req.site_name, chunk, page_url)
                            )
                            
                        cur_local.execute("INSERT OR REPLACE INTO web_crawl_cache (doc_id, url, content_hash, updated_at) VALUES (?, ?, ?, CURRENT_TIMESTAMP)", (doc_id, page_url, content_hash))
                        conn_local.commit()
                        conn_local.close()
                        
                    # 3. Prepare structure node
                    new_node = {
                        "title": pg['title'] or req.site_name,
                        "start_index": page_url,
                        "end_index": page_url,
                        "node_id": node_ids_for_page[0] if node_ids_for_page else str(uuid.uuid4()),
                        "text": chunks,
                        "summary": pg.get('summary', '요약 정보가 없습니다.')
                    }
                    return new_node

            tasks = [process_page(p) for p in pages]
            return await asyncio.gather(*tasks)

        if pages_to_summarize:
            document_status[doc_id]["progress"] = f"부분 수집됨 - LLM 요약 중 ({len(pages_to_summarize)}페이지)..."
            save_status()
            
            if not is_clear_existing:
                urls_to_update = {p.get("url") for p in pages_to_summarize if p.get("url")}
                current_structure = [n for n in current_structure if n.get("start_index") not in urls_to_update]
            
            import sys
            if sys.platform == 'win32':
                asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
            
            new_nodes = asyncio.run(summarize_and_index_all(pages_to_summarize))
            for node in new_nodes:
                if node:
                    current_structure.append(node)
            
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("SELECT COUNT(DISTINCT url) FROM web_crawl_cache WHERE doc_id = ?", (doc_id,))
        total_pages = cursor.fetchone()[0]
        conn.close()
        
        # 트리 구조 최종 저장
        tree_obj = {
            "doc_name": req.site_name,
            "structure": current_structure
        }
        with open(tree_path, "w", encoding="utf-8") as f:
            json.dump(tree_obj, f, indent=2, ensure_ascii=False)
        
        # 최종 성공 상태 업데이트
        document_status[doc_id]["status"] = "ready"
        is_cur_cancelled = doc_id in cancelled_jobs
        document_status[doc_id]["progress"] = "크롤링 중단됨 (부분 수집 완료)" if is_cancelled or is_cur_cancelled else "크롤링 완료"
        document_status[doc_id]["progress_percent"] = 100
        document_status[doc_id]["page_count"] = total_pages
        
        summary_accum = [n.get("summary", "") for n in current_structure[:3] if n.get("summary")]
        desc = " ".join(summary_accum)
        desc = desc[:200] + "..." if len(desc) > 200 else desc
        document_status[doc_id]["description"] = desc if desc else "웹 크롤링 문서입니다."
        document_status[doc_id]["doc_description"] = desc if desc else "웹 크롤링 문서입니다."
        
        document_status[doc_id]["updated_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        save_status()
        
    except Exception as e:
        if doc_id in document_status:
            document_status[doc_id]["status"] = "failed"
            document_status[doc_id]["progress"] = f"오류 발생: {str(e)}"
            save_status()


@app.post("/api/documents/crawl")
async def api_crawl_website(req: CrawlRequest, background_tasks: BackgroundTasks, current_user: dict = fastapi.Depends(get_current_user)):
    import uuid
    from app.rag_engine import DB_PATH, document_status, save_status
    import sqlite3
    
    try:
        doc_id = str(uuid.uuid4())[:12]
        final_name = f"[WEBSITE] {req.site_name}"
        owner_id = current_user["id"]
        
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        
        category_name = "General"
        visibility = "private"
        
        if req.folder_id:
            cursor.execute("SELECT name, visibility FROM categories WHERE id = ?", (req.folder_id,))
            row = cursor.fetchone()
            if row:
                category_name = row[0]
                visibility = row[1]
                
        upload_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # 초기 상태 저장용
        document_status[doc_id] = {
            "status": "pending",
            "progress": "크롤링 대기 중...",
            "progress_percent": 5,
            "name": final_name,
            "safe_filename": f"{req.site_name}.url",
            "upload_date": upload_date,
            "updated_at": upload_date,
            "page_count": "Unknown",
            "category": category_name,
            "visibility": visibility,
            "is_active": True,
            "owner_id": current_user["id"],
            "organization_id": current_user.get("organization_id"),
            "doc_description": f"크롤링된 웹사이트 문서입니다. URL: {req.url}",
            "file_path": f"[WEBSITE] {req.url}",
            "crawl_options": {
                "max_depth": req.max_depth,
                "max_pages": req.max_pages,
                "crawl_type": req.crawl_type,
                "strategy": req.strategy,
                "restrict_path": req.restrict_path,
                "use_ai_extraction": req.use_ai_extraction,
                "ai_extraction_prompt": req.ai_extraction_prompt,
                "clear_existing": req.clear_existing
            }
        }
        save_status()
        
        background_tasks.add_task(
            background_crawl_website, 
            doc_id, req, owner_id, current_user.get("organization_id"), category_name, visibility
        )
            
        return {"success": True, "doc_id": doc_id, "message": "크롤링 요청이 등록되었습니다."}
        
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"크롤링 환경 구성 중 오류: {str(e)}"})

# 크롤링된 사이트 대상 로그 조회 API
@app.get("/api/websites/{doc_id}/logs")
async def get_website_logs(doc_id: str, current_user: dict = fastapi.Depends(get_current_user)):
    from app.rag_engine import DB_PATH
    from app.rag_engine import document_status
    import sqlite3
    
    logs_dict = {}
    
    # 1. Load existing parsed logs from DB
    try:
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("""
            SELECT d.page_num, d.title, sum(length(d.text_content)), max(c.updated_at)
            FROM docs_fts d
            LEFT JOIN web_crawl_cache c ON d.doc_id = c.doc_id AND d.page_num = c.url
            WHERE d.doc_id = ? 
            GROUP BY d.page_num, d.title
        """, (doc_id,))
        rows = cursor.fetchall()
        conn.close()
        for r in rows:
            logs_dict[r[0]] = {"url": r[0], "title": r[1], "length": r[2], "updated_at": r[3] or "알 수 없음"}
    except Exception:
        pass
        
    # 2. Skip live_logs merging to ensure UI strictly displays fully completed pages page-by-page
    return {"logs": list(logs_dict.values())}

# 원문 조회용 API (웹 뷰어 등에 사용)
@app.get("/api/documents/{doc_id}/text")
async def get_document_raw_text(doc_id: str, current_user: dict = fastapi.Depends(get_current_user)):
    from app.rag_engine import DB_PATH
    import sqlite3
    try:
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("SELECT page_num, title, text_content FROM docs_fts WHERE doc_id = ? ORDER BY page_num", (doc_id,))
        rows = cursor.fetchall()
        print(f"DEBUG /api/documents/{doc_id}/text: Found {len(rows)} chunks in docs_fts")
        if len(rows) > 0:
            print(f"DEBUG First chunk sample: page_num={rows[0][0]}, title={rows[0][1]}, text_len={len(rows[0][2] or '')}")
        conn.close()
        
        # 청크들을 URL별로 이어 붙여서 하나의 섹션으로 구성합니다.
        content_by_url = {}
        title_by_url = {}
        for row in rows:
            url, title, text = row
            if url not in content_by_url:
                content_by_url[url] = []
                title_by_url[url] = title
            content_by_url[url].append(text)
            
        formatted_sections = []
        for url, texts in content_by_url.items():
            combined_text = "\n".join(texts)
            t = title_by_url[url]
            formatted_sections.append(f"## {t}\n**Source:** {url}\n\n{combined_text}\n")
            
        full_markdown = "\n---\n".join(formatted_sections)
        if not full_markdown:
            full_markdown = "*추출된 텍스트 컨텐츠가 없습니다.*"
            
        return {"text": full_markdown}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

class ScheduleUpdate(BaseModel):
    schedule: str

@app.patch("/api/websites/{doc_id}/schedule")
async def update_website_schedule(doc_id: str, payload: ScheduleUpdate, current_user: dict = fastapi.Depends(get_current_user)):
    from app.rag_engine import document_status, save_status
    print(f"PATCH SCHEDULE CALLED FOR doc_id={doc_id} with schedule={payload.schedule}")
    
    if doc_id in document_status:
        old_schedule = document_status[doc_id].get("auto_crawl_schedule", "disable")
        document_status[doc_id]["auto_crawl_schedule"] = payload.schedule
        
        # If toggling on or changing to a new active schedule, force an immediate run 
        if payload.schedule != "disable" and payload.schedule != old_schedule:
            document_status[doc_id]["last_auto_crawl_at"] = "2000-01-01 00:00:00"
        elif "last_auto_crawl_at" not in document_status[doc_id]:
            document_status[doc_id]["last_auto_crawl_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            
        save_status()
        print(f"SUCCESS: Saved {payload.schedule} for {doc_id}")
        return {"success": True}
    
    print(f"ERROR: doc_id {doc_id} NOT FOUND IN document_status")
    return JSONResponse(status_code=404, content={"error": "Website document not found."})

# 크롤링 멈춤 API
@app.post("/api/websites/{doc_id}/stop")
async def stop_crawl_website(doc_id: str, current_user: dict = fastapi.Depends(get_current_user)):
    from app.rag_engine import cancelled_jobs, document_status, save_status
    cancelled_jobs.add(doc_id)
    if doc_id in document_status and document_status[doc_id].get("status") in ["pending", "processing"]:
        document_status[doc_id]["progress"] = "중지 명령 수신 (종료 중...)"
        save_status()
    return {"success": True}

class RecrawlOptions(BaseModel):
    crawl_type: Optional[str] = "spa"
    use_ai_extraction: Optional[bool] = False
    ai_extraction_prompt: Optional[str] = None
    strategy: Optional[str] = "bfs"
    max_depth: Optional[int] = 3
    max_pages: Optional[int] = 50
    restrict_path: Optional[bool] = False
    clear_existing: Optional[bool] = True

# 재크롤링 API
@app.post("/api/websites/{doc_id}/recrawl")
async def recrawl_website(doc_id: str, options: RecrawlOptions, background_tasks: BackgroundTasks, current_user: dict = fastapi.Depends(get_current_user)):
    from app.rag_engine import document_status, DB_PATH, cancelled_jobs
    if doc_id not in document_status:
        return JSONResponse(status_code=404, content={"error": "문서를 찾을 수 없습니다."})
        
    doc_info = document_status[doc_id]
    
    # 방어 로직: 두 번 누르면 중복 실행되어 스레드 경합(꼬임)이 발생할 수 있음
    if doc_info.get("status") in ["processing", "crawling", "pending"]:
        return JSONResponse(status_code=400, content={"error": "현재 문서 작업이 처리(또는 안전하게 종료) 중입니다. 완전히 종료되어 상태가 변경될 때까지 몇 초만 기다린 후 재시도하세요."})
        
    cancelled_jobs.discard(doc_id)
    
    original_url = ""
    # URL 획득법: doc_description에서 가져오거나 file_path에서 획득
    fp = doc_info.get("file_path", "")
    if fp.startswith("[WEBSITE] "): original_url = fp.replace("[WEBSITE] ", "").strip()
    
    if not original_url: return JSONResponse(status_code=400, content={"error": "크롤링 설정 정보를 찾을 수 없습니다."})
        
    # 동일한 CrawlRequest를 가라로 생성
    req = CrawlRequest(
        url=original_url, 
        site_name=doc_info.get("name", "Unknown").replace("[WEBSITE] ", "").strip(), 
        max_depth=options.max_depth, 
        max_pages=options.max_pages,
        crawl_type=options.crawl_type,
        strategy=options.strategy,
        restrict_path=options.restrict_path,
        use_ai_extraction=options.use_ai_extraction,
        ai_extraction_prompt=options.ai_extraction_prompt,
        clear_existing=options.clear_existing
    )
    # 상태 즉각 초기화 (실시간 반영을 위함)
    doc_info["status"] = "processing"
    doc_info["progress"] = "재크롤링 시작 대기 중..."
    doc_info["progress_percent"] = 0
    doc_info["crawl_options"] = {
        "max_depth": options.max_depth,
        "max_pages": options.max_pages,
        "crawl_type": options.crawl_type,
        "strategy": options.strategy,
        "restrict_path": options.restrict_path,
        "use_ai_extraction": options.use_ai_extraction,
        "ai_extraction_prompt": options.ai_extraction_prompt,
        "clear_existing": options.clear_existing
    }
    if options.clear_existing:
        doc_info["live_logs"] = []
    
    background_tasks.add_task(
        background_crawl_website, 
        doc_id, req, current_user["id"], current_user.get("organization_id"), doc_info.get("category"), doc_info.get("visibility")
    )
    return {"success": True, "message": "재크롤링 요청이 큐에 담겼습니다."}



if __name__ == "__main__":

    import uvicorn

    uvicorn.run("app.main:app", host="0.0.0.0", port=8000, reload=True)






@app.delete("/api/chat/sessions")
async def delete_all_chat_sessions(current_user: dict = fastapi.Depends(get_current_user)):
    """Delete all chat sessions and messages for the current user."""
    try:
        conn = sqlite3.connect(str(DB_PATH), timeout=30)
        cursor = conn.cursor()
        cursor.execute("DELETE FROM chat_history WHERE user_id = ?", (current_user["id"],))
        cursor.execute("DELETE FROM chat_sessions WHERE user_id = ?", (current_user["id"],))
        conn.commit()
        conn.close()
        return {"message": "All sessions deleted"}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

